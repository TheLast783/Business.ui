"""Update presentation slides 3 to 6 for Warriors_ 3-Model Triad.pptx."""
import sys
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def update_presentation():
    prs = pptx.Presentation("Warriors_ 3-Model Triad.pptx")
    print(f"Loaded presentation with {len(prs.slides)} slides.")

    # -------------------------------------------------------------------------
    # SLIDE 3: Problem Statement
    # -------------------------------------------------------------------------
    slide3 = prs.slides[2]
    # Update title
    for shape in slide3.shapes:
        if "Title" in shape.name or (shape.has_text_frame and "problem statement" in shape.text_frame.text.lower()):
            shape.text_frame.text = "Describe the problem statement (200 words)"
            if shape.text_frame.paragraphs:
                p = shape.text_frame.paragraphs[0]
                p.font.name = "Arial"
                p.font.size = Pt(22)
                p.font.bold = True
                p.font.color.rgb = RGBColor(70, 20, 130) # Accenture deep purple/dark
        elif shape.has_text_frame and shape.name != "Title 3":
            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True
            
            p0 = tf.paragraphs[0]
            p0.text = "Modern enterprise dashboards excel at reporting WHAT happened, but fail completely at explaining WHY or WHAT NEXT. When a critical business KPI drops—e.g. gross revenue falls 12%—dashboards simply display a red alert. Leaders cannot distinguish between an internal operational breakdown (payment gateway timeout, ERP inventory backlog) and an external macroeconomic shock (port congestion, competitor flash promo)."
            p0.font.name = "Arial"
            p0.font.size = Pt(11)
            p0.font.color.rgb = RGBColor(40, 40, 40)
            p0.space_after = Pt(6)

            p1 = tf.add_paragraph()
            p1.text = "Consequently, root-cause interpretation falls onto human analysts who spend 4 to 7 days manually querying SQL databases, Jira tickets, and market feeds."
            p1.font.name = "Arial"
            p1.font.size = Pt(11)
            p1.font.color.rgb = RGBColor(40, 40, 40)
            p1.space_after = Pt(6)

            p2 = tf.add_paragraph()
            p2.text = "This human bottleneck introduces three critical enterprise risks:"
            p2.font.name = "Arial"
            p2.font.size = Pt(11)
            p2.font.bold = True
            p2.font.color.rgb = RGBColor(40, 40, 40)
            p2.space_after = Pt(4)

            risks = [
                ("Critical Latency: ", "By the time manual root cause is identified 4-7 days later, revenue loss becomes irreversible."),
                ("Alert Fatigue: ", "Teams waste high-value bandwidth investigating normal weekend and cyclical swings as false alarms."),
                ("Black-Box AI Guesswork & Leakage: ", "Generic LLMs hallucinate calculations, make arithmetic errors, and risk leaking sensitive margin/cost data across internal roles.")
            ]
            for title, desc in risks:
                pr = tf.add_paragraph()
                pr.space_after = Pt(4)
                pr.level = 0
                run1 = pr.add_run()
                run1.text = "• " + title
                run1.font.name = "Arial"
                run1.font.size = Pt(10.5)
                run1.font.bold = True
                run1.font.color.rgb = RGBColor(161, 0, 255) # Accenture Purple Accent
                
                run2 = pr.add_run()
                run2.text = desc
                run2.font.name = "Arial"
                run2.font.size = Pt(10.5)
                run2.font.color.rgb = RGBColor(50, 50, 50)

    # -------------------------------------------------------------------------
    # SLIDE 4: Proposed Solution - Architecture & Diagnostic Engine
    # -------------------------------------------------------------------------
    slide4 = prs.slides[3]
    for shape in slide4.shapes:
        if "Title" in shape.name or (shape.has_text_frame and "proposed solution" in shape.text_frame.text.lower()):
            shape.text_frame.text = "Proposed solution (200 words) — Diagnostic Engine"
            if shape.text_frame.paragraphs:
                p = shape.text_frame.paragraphs[0]
                p.font.name = "Arial"
                p.font.size = Pt(22)
                p.font.bold = True
                p.font.color.rgb = RGBColor(70, 20, 130)
        elif shape.has_text_frame and shape.name != "Title 3":
            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True

            p0 = tf.paragraphs[0]
            p0.text = "Our solution, BusinessIntelligence.ai (CausaMetric Engine), decouples deterministic mathematics from AI synthesis across a multi-tier cognitive intelligence loop:"
            p0.font.name = "Arial"
            p0.font.size = Pt(11)
            p0.font.color.rgb = RGBColor(40, 40, 40)
            p0.space_after = Pt(6)

            p1 = tf.add_paragraph()
            p1.text = "1. Deterministic Non-LLM Analytical Core (Zero Hallucinations):"
            p1.font.name = "Arial"
            p1.font.size = Pt(11)
            p1.font.bold = True
            p1.font.color.rgb = RGBColor(70, 20, 130)
            p1.space_after = Pt(3)

            core_bullets = [
                ("Statistical Process Control (SPC): ", "28-day Day-of-Week (DoW) normalized baseline with >2.5σ threshold mathematically filters seasonal noise, preventing alert fatigue."),
                ("Exact Shapley Causal Metric Tree: ", "Decomposes ΔRevenue = ΔSessions + ΔCVR + ΔAOV with strict 0.00 residual error (|ε| < 10⁻¹²) and 0 LLM math tokens.")
            ]
            for title, desc in core_bullets:
                pr = tf.add_paragraph()
                pr.space_after = Pt(3)
                run1 = pr.add_run()
                run1.text = "• " + title
                run1.font.name = "Arial"
                run1.font.size = Pt(10.5)
                run1.font.bold = True
                run1.font.color.rgb = RGBColor(0, 102, 204)
                run2 = pr.add_run()
                run2.text = desc
                run2.font.name = "Arial"
                run2.font.size = Pt(10.5)
                run2.font.color.rgb = RGBColor(50, 50, 50)

            p2 = tf.add_paragraph()
            p2.text = "2. Dual Diagnostic Intelligences (Parallel Execution):"
            p2.font.name = "Arial"
            p2.font.size = Pt(11)
            p2.font.bold = True
            p2.font.color.rgb = RGBColor(70, 20, 130)
            p2.space_before = Pt(4)
            p2.space_after = Pt(3)

            diag_bullets = [
                ("Model 1 (Private Internal Specialist): ", "Fine-tuned on 100k+ enterprise failure patterns, runs 100% on-premise to diagnose ERP backlogs and Jira tickets with zero data leakage."),
                ("Model 2 (Live Global Sentinel): ", "Connects to real-time market/news APIs to track external shocks (port congestion, competitor promos) invisible to internal databases.")
            ]
            for title, desc in diag_bullets:
                pr = tf.add_paragraph()
                pr.space_after = Pt(3)
                run1 = pr.add_run()
                run1.text = "• " + title
                run1.font.name = "Arial"
                run1.font.size = Pt(10.5)
                run1.font.bold = True
                run1.font.color.rgb = RGBColor(161, 0, 255)
                run2 = pr.add_run()
                run2.text = desc
                run2.font.name = "Arial"
                run2.font.size = Pt(10.5)
                run2.font.color.rgb = RGBColor(50, 50, 50)

    # -------------------------------------------------------------------------
    # SLIDE 5: Proposed Solution - Prescriptive Simulation & Governance
    # -------------------------------------------------------------------------
    slide5 = prs.slides[4]
    for shape in slide5.shapes:
        if "Title" in shape.name or (shape.has_text_frame and "proposed solution" in shape.text_frame.text.lower()):
            shape.text_frame.text = "Proposed solution (200 words) — Simulation & Action"
            if shape.text_frame.paragraphs:
                p = shape.text_frame.paragraphs[0]
                p.font.name = "Arial"
                p.font.size = Pt(22)
                p.font.bold = True
                p.font.color.rgb = RGBColor(70, 20, 130)
        elif shape.has_text_frame and "Title" not in shape.name:
            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True

            p1 = tf.paragraphs[0]
            p1.text = "3. Model 3: Prescriptive Action & 30/60/90-Day Trajectory ROI Simulator:"
            p1.font.name = "Arial"
            p1.font.size = Pt(11)
            p1.font.bold = True
            p1.font.color.rgb = RGBColor(70, 20, 130)
            p1.space_after = Pt(3)

            m3_bullets = [
                ("Attribution Reconciliation: ", "Synthesizes Model 1 & 2 findings into exact causal shares (e.g. 70% port strike + 30% warehouse backlog)."),
                ("Pre-Action Outcome Simulation: ", "Projects revenue over 30/60/90 days for Status Quo, Prescribed Recovery, and Constrained paths with forecasted ROI multipliers before spending capital."),
                ("Human Mind Mixing (Executive Co-Pilot): ", "Dynamic budget sliders ($45k cap) and policy overrides instantly recalculate recovery trajectories in real time.")
            ]
            for title, desc in m3_bullets:
                pr = tf.add_paragraph()
                pr.space_after = Pt(3)
                run1 = pr.add_run()
                run1.text = "• " + title
                run1.font.name = "Arial"
                run1.font.size = Pt(10.5)
                run1.font.bold = True
                run1.font.color.rgb = RGBColor(0, 102, 204)
                run2 = pr.add_run()
                run2.text = desc
                run2.font.name = "Arial"
                run2.font.size = Pt(10.5)
                run2.font.color.rgb = RGBColor(50, 50, 50)

            p2 = tf.add_paragraph()
            p2.text = "4. Governed Security, Explicit Abstention & Persona Tailoring:"
            p2.font.name = "Arial"
            p2.font.size = Pt(11)
            p2.font.bold = True
            p2.font.color.rgb = RGBColor(70, 20, 130)
            p2.space_before = Pt(4)
            p2.space_after = Pt(3)

            gov_bullets = [
                ("Explicit Abstention Protocol: ", "When evidence confidence margin is <25%, the engine halts high-capital action and prescribes 2 low-cost canary validation tests ($120-$350)."),
                ("Role-Based Access Control (RBAC): ", "Masks sensitive unit COGS and margins for Operations Analysts while providing full financial transparency to Executives.")
            ]
            for title, desc in gov_bullets:
                pr = tf.add_paragraph()
                pr.space_after = Pt(3)
                run1 = pr.add_run()
                run1.text = "• " + title
                run1.font.name = "Arial"
                run1.font.size = Pt(10.5)
                run1.font.bold = True
                run1.font.color.rgb = RGBColor(161, 0, 255)
                run2 = pr.add_run()
                run2.text = desc
                run2.font.name = "Arial"
                run2.font.size = Pt(10.5)
                run2.font.color.rgb = RGBColor(50, 50, 50)

    # -------------------------------------------------------------------------
    # SLIDE 6: Thank You & Prototype Verification Summary
    # -------------------------------------------------------------------------
    slide6 = prs.slides[5]
    for shape in slide6.shapes:
        if "Title" in shape.name or (shape.has_text_frame and "thank you" in shape.text_frame.text.lower()):
            shape.text_frame.text = "Thank you — BusinessIntelligence.ai"
            if shape.text_frame.paragraphs:
                p = shape.text_frame.paragraphs[0]
                p.font.name = "Arial"
                p.font.size = Pt(22)
                p.font.bold = True
                p.font.color.rgb = RGBColor(70, 20, 130)

    # Add a clean summary text box on slide 6 if not present
    slide6_textbox = None
    for shape in slide6.shapes:
        if shape.has_text_frame and "Title" not in shape.name:
            slide6_textbox = shape
            break
    
    if not slide6_textbox:
        left = Inches(1.0)
        top = Inches(2.2)
        width = Inches(11.3)
        height = Inches(4.5)
        slide6_textbox = slide6.shapes.add_textbox(left, top, width, height)

    tf6 = slide6_textbox.text_frame
    tf6.clear()
    tf6.word_wrap = True

    p0 = tf6.paragraphs[0]
    p0.text = "Delivered Working Prototype & Verification Highlights:"
    p0.font.name = "Arial"
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = RGBColor(70, 20, 130)
    p0.space_after = Pt(8)

    summary_items = [
        ("Interactive Decision Workspace: ", "Complete Python + Streamlit decision application with scenario selector, persona switcher, Plotly SPC control charts, Shapley waterfalls, and 30/60/90-day trajectory simulations."),
        ("142 Automated Passing Tests: ", "93/93 headless scenario acceptance tests + 49/49 property tests with 10,000 Monte Carlo trials verifying zero residual error (|ε| < 10⁻¹²) and 0 LLM math tokens."),
        ("All 4 Mandatory Scenarios Implemented: ", "1. Multi-factor compound drop (70/30) • 2. Low-confidence ambiguity with canary tests • 3. Sparse-history cold start launch • 4. Role-based data masking."),
        ("Pluggable & Secure Architecture: ", "Seamlessly toggles between live API providers (Gemini, OpenAI, Ollama) and an instant deterministic offline fallback (<5ms latency, $0.00 cost).")
    ]

    for title, desc in summary_items:
        pr = tf6.add_paragraph()
        pr.space_after = Pt(6)
        run1 = pr.add_run()
        run1.text = "• " + title
        run1.font.name = "Arial"
        run1.font.size = Pt(11)
        run1.font.bold = True
        run1.font.color.rgb = RGBColor(161, 0, 255)
        run2 = pr.add_run()
        run2.text = desc
        run2.font.name = "Arial"
        run2.font.size = Pt(11)
        run2.font.color.rgb = RGBColor(50, 50, 50)

    p_team = tf6.add_paragraph()
    p_team.space_before = Pt(12)
    p_team.text = "Team: Warriors (IIT Patna) — Harsh Singh Baghel, Rathod Rudra, Dhruv Maheshwari"
    p_team.font.name = "Arial"
    p_team.font.size = Pt(11)
    p_team.font.italic = True
    p_team.font.color.rgb = RGBColor(100, 100, 100)

    # Save presentation
    prs.save("Warriors_ 3-Model Triad.pptx")
    print("Successfully updated Warriors_ 3-Model Triad.pptx!")

if __name__ == "__main__":
    update_presentation()
