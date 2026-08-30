"""
Clean up Warriors_BusinessIntelligence_Accenture_2026_FONT_FIXED.pptx:
Removes accidental middle duplicate slide 6 while strictly preserving every single color, shape, font, and layout.
"""
import os
from pptx import Presentation

def clean_deck(input_path: str, output_path: str):
    prs = Presentation(input_path)
    
    # Check if slide 6 is the duplicate "Thank you" slide
    # Slide index 5 in 0-indexed list
    s6 = prs.slides[5]
    s6_texts = [sh.text.strip() for sh in s6.shapes if hasattr(sh, 'text') and sh.text.strip()]
    
    if len(s6_texts) == 1 and "Thank you" in s6_texts[0] and len(prs.slides) > 6:
        print(f"Removing middle duplicate slide 6: {s6_texts}")
        rId = prs.slides._sldIdLst[5].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[5]

    prs.save(output_path)
    print(f"Cleaned presentation saved with {len(prs.slides)} slides to: {output_path}")

if __name__ == "__main__":
    in_path = "Warriors_BusinessIntelligence_Accenture_2026_FONT_FIXED.pptx"
    out_path = "Warriors_BusinessIntelligence_Accenture_2026_FINAL.pptx"
    clean_deck(in_path, out_path)
    
    # Also overwrite original if not locked
    try:
        clean_deck(in_path, in_path)
    except PermissionError:
        print("Note: Warriors_BusinessIntelligence_Accenture_2026_FONT_FIXED.pptx is currently open in PowerPoint.")
