"""
Generate the official Executive Presentation (.pptx) for BusinessIntelligence.ai
Design: 16:9 Widescreen, Enterprise Navy/Slate palette, Bigger Fonts, Medium Readable Content.
"""
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck(output_path: str):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette
    COLOR_BG = RGBColor(9, 14, 23)          # #090E17
    COLOR_CARD = RGBColor(15, 23, 42)       # #0F172A
    COLOR_BORDER = RGBColor(30, 41, 59)     # #1E293B
    COLOR_PRIMARY = RGBColor(2, 132, 199)   # #0284C7
    COLOR_ACCENT = RGBColor(56, 189, 248)   # #38BDF8
    COLOR_TEXT = RGBColor(248, 250, 252)    # #F8FAFC
    COLOR_MUTED = RGBColor(148, 163, 184)   # #94A3B8
    COLOR_GREEN = RGBColor(16, 185, 129)    # #10B981
    COLOR_AMBER = RGBColor(245, 158, 11)    # #F59E0B

    blank_layout = prs.slide_layouts[6]

    def set_slide_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.color.rgb = COLOR_BG
        return bg

    def add_header(slide, title_text: str, category_text: str = "BUSINESSINTELLIGENCE.AI · ROUND 2"):
        # Category Tracker
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_ACCENT

        # Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.7), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT

    def add_card(slide, left, top, width, height, bg_color=COLOR_CARD, border_color=COLOR_BORDER):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
        return shape

    # ==========================================
    # SLIDE 1: TITLE SLIDE
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s1)

    # Accent bar
    bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(2.2), Inches(0.12), Inches(3.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.color.rgb = COLOR_PRIMARY

    t_box = s1.shapes.add_textbox(Inches(1.6), Inches(2.1), Inches(10.5), Inches(3.4))
    tf1 = t_box.text_frame
    tf1.word_wrap = True

    p1 = tf1.paragraphs[0]
    p1.text = "BusinessIntelligence.ai"
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TEXT

    p2 = tf1.add_paragraph()
    p2.text = "KPI Root Cause Diagnosis & Prescriptive Action Engine"
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_ACCENT
    p2.space_before = Pt(12)

    p3 = tf1.add_paragraph()
    p3.text = "Deterministic Non-LLM Math Core  ·  3-Model AI Synthesis  ·  Governed Semantic Contract"
    p3.font.size = Pt(15)
    p3.font.color.rgb = COLOR_MUTED
    p3.space_before = Pt(16)

    # Bottom Metadata Box
    add_card(s1, Inches(1.6), Inches(5.8), Inches(10.0), Inches(0.9))
    meta_box = s1.shapes.add_textbox(Inches(1.8), Inches(5.85), Inches(9.6), Inches(0.8))
    tf_meta = meta_box.text_frame
    p_meta = tf_meta.paragraphs[0]
    p_meta.text = "Accenture AI Hackathon — Round 2 Solution  |  Multi-Source Intelligence-to-Action Engine"
    p_meta.font.size = Pt(13)
    p_meta.font.bold = True
    p_meta.font.color.rgb = COLOR_TEXT

    # ==========================================
    # SLIDE 2: EXECUTIVE SUMMARY & PROBLEM STATEMENT
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s2)
    add_header(s2, "Executive Problem: The Telemetry-to-Action Gap")

    # Card 1: The Challenge
    add_card(s2, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0))
    tb_c1 = s2.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.0), Inches(4.5))
    tf_c1 = tb_c1.text_frame
    tf_c1.word_wrap = True

    p = tf_c1.paragraphs[0]
    p.text = "The Operational Dilemma"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_AMBER

    points_c1 = [
        ("Noisy Anomaly Alerts: ", "Static dashboards flood leaders with false alarms, failing to isolate seasonal Day-of-Week swings from genuine operational breaks."),
        ("Multi-Source Data Silos: ", "ERP SQL ledgers (daily), web session telemetry (hourly), and Jira tickets (weekly) are disconnected."),
        ("The LLM Hallucination Risk: ", "Pure LLM analytics fabricate mathematical totals, violate financial identities, and propose generic, ungrounded actions.")
    ]
    for bold_t, norm_t in points_c1:
        p = tf_c1.add_paragraph()
        p.space_before = Pt(14)
        run1 = p.add_run()
        run1.text = bold_t
        run1.font.bold = True
        run1.font.size = Pt(14)
        run1.font.color.rgb = COLOR_TEXT
        run2 = p.add_run()
        run2.text = norm_t
        run2.font.size = Pt(13)
        run2.font.color.rgb = COLOR_MUTED

    # Card 2: The Solution
    add_card(s2, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb_c2 = s2.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf_c2 = tb_c2.text_frame
    tf_c2.word_wrap = True

    p = tf_c2.paragraphs[0]
    p.text = "Our Architectural Solution"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN

    points_c2 = [
        ("Deterministic Non-LLM Math Core: ", "Statistical Process Control (28-day DoW baseline) + exact 3-Factor Shapley decomposition with guaranteed $0.00 residual error."),
        ("3-Model AI Diagnostic Triad: ", "Internal diagnostic (Model 1) + Live macro sentinel (Model 2) + Prescriptive ROI trajectory simulator (Model 3)."),
        ("Governed Safe Action Engine: ", "Explicit abstention protocol under ambiguity ($<25\\%$ confidence delta) prescribing low-cost canary validation tests.")
    ]
    for bold_t, norm_t in points_c2:
        p = tf_c2.add_paragraph()
        p.space_before = Pt(14)
        run1 = p.add_run()
        run1.text = bold_t
        run1.font.bold = True
        run1.font.size = Pt(14)
        run1.font.color.rgb = COLOR_TEXT
        run2 = p.add_run()
        run2.text = norm_t
        run2.font.size = Pt(13)
        run2.font.color.rgb = COLOR_MUTED

    # ==========================================
    # SLIDE 3: GOVERNED DATA LAYER & SEMANTIC CONTRACT
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s3)
    add_header(s3, "Governed Semantic Contract & Heterogeneous Ingestion")

    # 3 Column Cards for 3 Data Sources
    sources = [
        ("Daily ERP Sales Transactions", "Daily Grain (Structured SQL)", "Contains transaction orders, SKU categories, quantities, unit prices, and confidential COGS.", "1,420 rows  ·  24h SLA", Inches(0.8)),
        ("Hourly Web Analytics Stream", "Hourly Grain (Event Stream)", "Captures visitor sessions, funnel drop-offs, bounce rates, and checkout initiation events.", "696 hourly bins  ·  1h SLA", Inches(4.8)),
        ("Weekly Support & Jira Logs", "Weekly Grain (Semi-Structured)", "Customer friction tickets, warehouse error logs, and shipping delay incident taxonomy.", "18 incidents  ·  168h SLA", Inches(8.8))
    ]

    for title, grain, desc, sla, left_pos in sources:
        add_card(s3, left_pos, Inches(1.8), Inches(3.7), Inches(3.4))
        tb = s3.shapes.add_textbox(left_pos + Inches(0.2), Inches(2.0), Inches(3.3), Inches(3.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT

        p_sub = tf.add_paragraph()
        p_sub.text = grain
        p_sub.font.size = Pt(12)
        p_sub.font.color.rgb = COLOR_GREEN
        p_sub.space_before = Pt(4)

        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = COLOR_MUTED
        p_desc.space_before = Pt(10)

        p_sla = tf.add_paragraph()
        p_sla.text = sla
        p_sla.font.size = Pt(12)
        p_sla.font.bold = True
        p_sla.font.color.rgb = COLOR_TEXT
        p_sla.space_before = Pt(12)

    # Bottom Full Width Card: The 4 Connected KPIs
    add_card(s3, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.4))
    tb_kpi = s3.shapes.add_textbox(Inches(1.1), Inches(5.6), Inches(11.1), Inches(1.1))
    tf_kpi = tb_kpi.text_frame
    tf_kpi.word_wrap = True
    p = tf_kpi.paragraphs[0]
    p.text = "Governed 4-KPI Connected Equation (Strict Mathematical Lineage)"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT

    p_eq = tf_kpi.add_paragraph()
    p_eq.text = "Gross Revenue (R) = Sessions (S) × Conversion Rate (CR) × Average Order Value (AOV)"
    p_eq.font.size = Pt(16)
    p_eq.font.bold = True
    p_eq.font.color.rgb = COLOR_ACCENT
    p_eq.space_before = Pt(4)

    # ==========================================
    # SLIDE 4: DETERMINISTIC NON-LLM MATHEMATICAL CORE
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s4)
    add_header(s4, "Deterministic Non-LLM Analytical Core")

    # Left: SPC
    add_card(s4, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0))
    tb_spc = s4.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.0), Inches(4.5))
    tf_spc = tb_spc.text_frame
    tf_spc.word_wrap = True
    p = tf_spc.paragraphs[0]
    p.text = "1. Statistical Process Control (SPC)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    spc_bullets = [
        ("Signal vs. Noise Filtering: ", "Applies 28-day rolling Day-of-Week seasonality index to prevent false alarms on regular weekend volume dips."),
        ("Dynamic Control Limits: ", "Computes UCL and LCL at $\\pm 2.5\\sigma$. Identifies Day 29 anomaly ($z = -20.75\\sigma$) as a true operational failure."),
        ("Zero LLM Tokens: ", "Executed purely in vectorized NumPy/Pandas in $<5\\text{ms}$ with zero API cost.")
    ]
    for b_title, b_desc in spc_bullets:
        p = tf_spc.add_paragraph()
        p.space_before = Pt(12)
        r1 = p.add_run()
        r1.text = b_title
        r1.font.bold = True
        r1.font.size = Pt(13)
        r1.font.color.rgb = COLOR_TEXT
        r2 = p.add_run()
        r2.text = b_desc
        r2.font.size = Pt(12)
        r2.font.color.rgb = COLOR_MUTED

    # Right: Shapley Decomposition
    add_card(s4, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb_tree = s4.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf_tree = tb_tree.text_frame
    tf_tree.word_wrap = True
    p = tf_tree.paragraphs[0]
    p.text = "2. Exact 3-Factor Shapley Decomposition"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN

    tree_bullets = [
        ("Exact Closed-Form Math: ", "Decomposes $-\\$65,600$ revenue drop across all factor permutation orderings (Sessions, CVR, AOV)."),
        ("Attribution Breakdown: ", "• Volume Impact: -$35,166 (-53.6%)\n• Conversion Rate: -$22,667 (-34.6%)\n• Basket Size (AOV): -$7,767 (-11.8%)"),
        ("Zero Residual Guarantee: ", "Sum of decomposed factors equals total revenue delta exactly: Residual $\\epsilon = \\$0.00$.")
    ]
    for b_title, b_desc in tree_bullets:
        p = tf_tree.add_paragraph()
        p.space_before = Pt(12)
        r1 = p.add_run()
        r1.text = b_title
        r1.font.bold = True
        r1.font.size = Pt(13)
        r1.font.color.rgb = COLOR_TEXT
        r2 = p.add_run()
        r2.text = b_desc
        r2.font.size = Pt(12)
        r2.font.color.rgb = COLOR_MUTED

    # ==========================================
    # SLIDE 5: 3-MODEL AI SYNTHESIS ARCHITECTURE
    # ==========================================
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s5)
    add_header(s5, "3-Model Cognitive AI Synthesis Triad")

    models = [
        ("Model 1: Diagnostic Engine", "Internal Operational Causes", "Analyzes internal ERP staging queues and Jira tickets (JIRA-4819). Isolates a 30% internal backlog at warehouse WH-WEST-01.", "30% Attribution Share", Inches(0.8), COLOR_PRIMARY),
        ("Model 2: Macro Sentinel", "External Live Market Feeds", "Ingests maritime port congestion feeds (MACRO-PORT-01). Quantifies a 70% external macro shock from West Coast labor slowdowns.", "70% Attribution Share", Inches(4.8), COLOR_AMBER),
        ("Model 3: Prescriptive Action", "ROI & Trajectory Simulator", "Blends Model 1 & 2 attributions, factors in dynamic executive constraints, and simulates 30/60/90-day recovery trajectories.", "4.2x Projected ROI", Inches(8.8), COLOR_GREEN)
    ]

    for title, subtitle, desc, metric, left_pos, border_c in models:
        add_card(s5, left_pos, Inches(1.8), Inches(3.7), Inches(5.0), border_color=border_c)
        tb = s5.shapes.add_textbox(left_pos + Inches(0.2), Inches(2.0), Inches(3.3), Inches(4.5))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT

        p_sub = tf.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.size = Pt(12)
        p_sub.font.color.rgb = border_c
        p_sub.space_before = Pt(4)

        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = COLOR_MUTED
        p_desc.space_before = Pt(14)

        p_met = tf.add_paragraph()
        p_met.text = metric
        p_met.font.size = Pt(14)
        p_met.font.bold = True
        p_met.font.color.rgb = border_c
        p_met.space_before = Pt(20)

    # ==========================================
    # SLIDE 6: SAFETY, CONFIDENCE & EXPLICIT ABSTENTION
    # ==========================================
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s6)
    add_header(s6, "AI Safety: Explicit Abstention & Canary Tests (Scenario 2)")

    add_card(s6, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0))
    tb_ab1 = s6.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.0), Inches(4.5))
    tf_ab1 = tb_ab1.text_frame
    tf_ab1.word_wrap = True

    p = tf_ab1.paragraphs[0]
    p.text = "When & Why We Abstain"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_AMBER

    ab_points = [
        ("The Capital Allocation Hazard: ", "Under ambiguous signals, forcing a single generative LLM prediction risks misallocating hundreds of thousands in capital."),
        ("Mathematical Abstention Gate: ", "Triggered automatically when the confidence margin between competing root causes is $<25\\%$ (observed: $16\\%$ delta)."),
        ("Dual Competing Hypotheses: ", "• Hypothesis 1 (60%): Payment Gateway 504 Timeout on iOS Checkout.\n• Hypothesis 2 (40%): Competitor Viral TikTok Flash Sale.")
    ]
    for b_title, b_desc in ab_points:
        p = tf_ab1.add_paragraph()
        p.space_before = Pt(12)
        r1 = p.add_run()
        r1.text = b_title
        r1.font.bold = True
        r1.font.size = Pt(13)
        r1.font.color.rgb = COLOR_TEXT
        r2 = p.add_run()
        r2.text = b_desc
        r2.font.size = Pt(12)
        r2.font.color.rgb = COLOR_MUTED

    # Right: Canary Tests
    add_card(s6, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb_ab2 = s6.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf_ab2 = tb_ab2.text_frame
    tf_ab2.word_wrap = True

    p = tf_ab2.paragraphs[0]
    p.text = "Prescribed Low-Cost Discovery Tests"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN

    canary_tests = [
        ("CANARY TEST 1: Synthetic iOS Checkout Probe", "Simulate 500 programmatic checkout requests through iOS gateway corridor.\nCost: $150  |  Duration: 2.0 hours  |  Decision Gate: Timeout Rate > 5%"),
        ("CANARY TEST 2: Competitive Ad & Social Scraping Run", "Ingest competitor TikTok and Instagram ad feeds across ASEAN.\nCost: $300  |  Duration: 4.0 hours  |  Decision Gate: Competitor Share of Voice Spike > 25%")
    ]
    for title, body in canary_tests:
        p = tf_ab2.add_paragraph()
        p.space_before = Pt(14)
        r1 = p.add_run()
        r1.text = title + "\n"
        r1.font.bold = True
        r1.font.size = Pt(13)
        r1.font.color.rgb = COLOR_ACCENT
        r2 = p.add_run()
        r2.text = body
        r2.font.size = Pt(12)
        r2.font.color.rgb = COLOR_MUTED

    # ==========================================
    # SLIDE 7: PERSONA DIVERGENCE & RBAC SECURITY
    # ==========================================
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s7)
    add_header(s7, "Persona Prescriptive Storytelling & RBAC Security")

    # Executive Persona
    add_card(s7, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0))
    tb_exec = s7.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.0), Inches(4.5))
    tf_exec = tb_exec.text_frame
    tf_exec.word_wrap = True
    p = tf_exec.paragraphs[0]
    p.text = "Executive / VP Persona"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    exec_items = [
        ("Tailored Decision Lens: ", "High-level strategic trade-offs, financial risk exposure, capital ROI, and supply chain corridor reallocation."),
        ("Prescribed Interventions: ", "• Authorize 35% maritime freight reroute to Seattle port corridors.\n• Authorize $45,000 expedited drayage budget to clear high-margin inventory."),
        ("Unmasked Financial Data: ", "Full numeric transparency over unit COGS, gross margin amounts, and margin percentages.")
    ]
    for b_title, b_desc in exec_items:
        p = tf_exec.add_paragraph()
        p.space_before = Pt(12)
        r1 = p.add_run()
        r1.text = b_title
        r1.font.bold = True
        r1.font.size = Pt(13)
        r1.font.color.rgb = COLOR_TEXT
        r2 = p.add_run()
        r2.text = b_desc
        r2.font.size = Pt(12)
        r2.font.color.rgb = COLOR_MUTED

    # Operations Analyst Persona
    add_card(s7, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb_ops = s7.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf_ops = tb_ops.text_frame
    tf_ops.word_wrap = True
    p = tf_ops.paragraphs[0]
    p.text = "Operations / Category Manager"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    ops_items = [
        ("Tailored Decision Lens: ", "Tactical SOPs, warehouse triage playbooks, system container codes, and ticket resolution."),
        ("Prescribed Interventions: ", "• Execute automated inventory reallocation for SKU-ELEC-409.\n• Triage Jira incident JIRA-4819 at Long Beach Berth 4.\n• Trigger customer delivery ETA updates to enterprise accounts."),
        ("Enforced RBAC Masking: ", "Unit COGS and Gross Margins are strictly redacted with [RESTRICTED-EXEC].")
    ]
    for b_title, b_desc in ops_items:
        p = tf_ops.add_paragraph()
        p.space_before = Pt(12)
        r1 = p.add_run()
        r1.text = b_title
        r1.font.bold = True
        r1.font.size = Pt(13)
        r1.font.color.rgb = COLOR_TEXT
        r2 = p.add_run()
        r2.text = b_desc
        r2.font.size = Pt(12)
        r2.font.color.rgb = COLOR_MUTED

    # ==========================================
    # SLIDE 8: HUMAN FEEDBACK & TELEMETRY AUDIT
    # ==========================================
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s8)
    add_header(s8, "Human-in-the-Loop Feedback & Telemetry Audit")

    # Left: Feedback Loop
    add_card(s8, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0))
    tb_fb = s8.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.0), Inches(4.5))
    tf_fb = tb_fb.text_frame
    tf_fb.word_wrap = True
    p = tf_fb.paragraphs[0]
    p.text = "Continuous Learning Loop"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    fb_items = [
        ("Governed Human Corrections: ", "Analysts submit star ratings, categorical tags, and operational driver corrections."),
        ("Cognitive Interpretation: ", "AI classifies feedback type, identifies target layer, and generates learning signals."),
        ("Bayesian Calibration Weights: ", "Updates operational driver ranking weights over time."),
        ("Quantitative Lock Principle: ", "Closed-form formulas, variance math, and SPC thresholds remain strictly immutable.")
    ]
    for b_title, b_desc in fb_items:
        p = tf_fb.add_paragraph()
        p.space_before = Pt(10)
        r1 = p.add_run()
        r1.text = b_title
        r1.font.bold = True
        r1.font.size = Pt(13)
        r1.font.color.rgb = COLOR_TEXT
        r2 = p.add_run()
        r2.text = b_desc
        r2.font.size = Pt(12)
        r2.font.color.rgb = COLOR_MUTED

    # Right: Telemetry Box
    add_card(s8, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb_tel = s8.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf_tel = tb_tel.text_frame
    tf_tel.word_wrap = True
    p = tf_tel.paragraphs[0]
    p.text = "Runtime Telemetry & Cost Accounting"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN

    tel_items = [
        ("Total Pipeline Latency: ", "< 2,000 ms (Ingest 1.5ms | SPC & Tree Math 4.2ms | AI Synthesis 1,600ms)."),
        ("Deterministic Math Tokens: ", "0 Tokens consumed (100% Non-LLM Closed-Form Python/NumPy)."),
        ("AI Narrative Tokens: ", "630 Tokens (450 prompt + 180 completion tokens)."),
        ("Cost Transparency: ", "$0.001035 estimated cost per execution (with $0.00 offline fallback).")
    ]
    for b_title, b_desc in tel_items:
        p = tf_tel.add_paragraph()
        p.space_before = Pt(10)
        r1 = p.add_run()
        r1.text = b_title
        r1.font.bold = True
        r1.font.size = Pt(13)
        r1.font.color.rgb = COLOR_TEXT
        r2 = p.add_run()
        r2.text = b_desc
        r2.font.size = Pt(12)
        r2.font.color.rgb = COLOR_MUTED

    # ==========================================
    # SLIDE 9: 4 BENCHMARK SCENARIOS MATRIX
    # ==========================================
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s9)
    add_header(s9, "Complete 4-Scenario Validation Matrix")

    scenario_cards = [
        ("Scenario 1: Multi-Factor Disruption", "Compound Drivers", "70% External West Coast port strike + 30% local warehouse backlog. 4.2x projected ROI recovery.", Inches(0.8), Inches(1.8)),
        ("Scenario 2: Low-Confidence Ambiguity", "Explicit Abstention", "Conflicting evidence (60% vs 40%). Engine abstains and prescribes $150 canary validation experiments.", Inches(6.8), Inches(1.8)),
        ("Scenario 3: Sparse History / Launch", "Cold Start Prior", "New SKU launch with N=7 days history. Blends Bayesian category benchmark ($5,000 base) and staged pilots.", Inches(0.8), Inches(4.5)),
        ("Scenario 4: Role-Based Security", "RBAC Entitlements", "Enforces column/row masking on unit COGS and margins for Analysts while showing system IDs.", Inches(6.8), Inches(4.5))
    ]

    for title, subtitle, desc, left_pos, top_pos in scenario_cards:
        add_card(s9, left_pos, top_pos, Inches(5.7), Inches(2.4))
        tb = s9.shapes.add_textbox(left_pos + Inches(0.2), top_pos + Inches(0.15), Inches(5.3), Inches(2.1))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT

        p_sub = tf.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.size = Pt(12)
        p_sub.font.bold = True
        p_sub.font.color.rgb = COLOR_GREEN
        p_sub.space_before = Pt(2)

        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = COLOR_MUTED
        p_desc.space_before = Pt(6)

    # ==========================================
    # SLIDE 10: ARCHITECTURE & CONCLUSION
    # ==========================================
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s10)
    add_header(s10, "Summary: Production-Grade KPI Intelligence Engine")

    add_card(s10, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tb_end = s10.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(11.0), Inches(4.4))
    tf_end = tb_end.text_frame
    tf_end.word_wrap = True

    p = tf_end.paragraphs[0]
    p.text = "Key Architectural Differentiators"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT

    diffs = [
        ("1. Mathematical Rigor First: ", "Zero reliance on LLM calculation. SPC anomaly detection and Shapley causal attribution are 100% deterministic with $0.00 residual."),
        ("2. Multi-Model Context Synthesis: ", "Combines private ERP/Jira operational context with live macro market shocks into actionable, persona-tailored playbooks."),
        ("3. Enterprise Guardrails: ", "Explicit abstention under ambiguity, dynamic RBAC masking, and governed human feedback calibration."),
        ("4. Verified & Ready: ", "93 / 93 automated tests passed. 8 / 8 scenario-persona combinations validated. Deployed live on Streamlit Cloud.")
    ]
    for b_title, b_desc in diffs:
        p = tf_end.add_paragraph()
        p.space_before = Pt(14)
        r1 = p.add_run()
        r1.text = b_title
        r1.font.bold = True
        r1.font.size = Pt(15)
        r1.font.color.rgb = COLOR_ACCENT
        r2 = p.add_run()
        r2.text = b_desc
        r2.font.size = Pt(14)
        r2.font.color.rgb = COLOR_MUTED

    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    out_file = os.path.join(os.path.dirname(os.path.abspath(__file__)), "BusinessIntelligence_Executive_Deck.pptx")
    create_deck(out_file)
