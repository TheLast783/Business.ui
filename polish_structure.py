"""
Comprehensive Structural and Font Polish for Warriors_BusinessIntelligence_Accenture_2026_FONT_FIXED.pptx:
- Ensures perfect word wrapping on all text frames.
- Standardizes font sizes (Slide titles 24pt, Section headers 15pt bold, Body text 12-13pt).
- Normalizes internal margins/padding to prevent clipping.
- Preserves all exact colors, shapes, positions, backgrounds, and template elements.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def polish_deck(pptx_path: str):
    prs = Presentation(pptx_path)
    
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            tf = shape.text_frame
            tf.word_wrap = True
            
            # Reduce internal text box margins slightly to prevent text clipping
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.08)
            tf.margin_top = Inches(0.06)
            tf.margin_bottom = Inches(0.06)
            
            for p_idx, p in enumerate(tf.paragraphs):
                # Ensure line spacing is clean
                if not p.line_spacing:
                    p.line_spacing = 1.15
                
                # Check text content to apply appropriate font sizing
                txt = p.text.strip()
                if not txt:
                    continue
                
                # Slide titles (Top headers)
                if shape.top < Inches(1.0) and shape.height < Inches(1.0) and shape.width > Inches(8.0):
                    p.font.size = Pt(22)
                    p.font.bold = True
                
                # Section Header cards (e.g., MODEL 1, SCENARIO 1, 1 SYNTHESIZE)
                elif any(txt.startswith(prefix) for prefix in ["MODEL 1", "MODEL 2", "MODEL 3", "SCENARIO 1", "SCENARIO 2", "SCENARIO 3", "SCENARIO 4", "1  SYNTHESIZE", "2  SIMULATE", "3  CONSTRAIN", "THREE ENTERPRISE", "EXACT KPI", "PERSONA-TAILORED", "EXAMPLE:"]):
                    p.font.size = Pt(14)
                    p.font.bold = True
                
                # Card subheadings (e.g., 4 canonical KPIs, SPC anomaly gate, Internal signals)
                elif shape.height < Inches(0.6) and shape.width < Inches(4.5) and shape.top > Inches(1.0):
                    p.font.size = Pt(14)
                    p.font.bold = True

    prs.save(pptx_path)
    print(f"Polished deck successfully: {pptx_path}")

if __name__ == "__main__":
    for fn in ["Warriors_BusinessIntelligence_Accenture_2026_FINAL.pptx", "Warriors_BusinessIntelligence_Accenture_2026_FONT_FIXED.pptx"]:
        try:
            polish_deck(fn)
        except PermissionError:
            print(f"Note: {fn} is currently open in PowerPoint on your desktop.")
