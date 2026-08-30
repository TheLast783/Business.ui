"""
Fix all formatting, overlapping placeholders, and ghost text in Warriors_BusinessIntelligence_Accenture_2026_FONT_FIXED.pptx
"""
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def fix_presentation(input_file: str, output_file: str):
    prs = Presentation(input_file)
    
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_PURPLE = RGBColor(161, 0, 255)       # Accenture Purple #A100FF
    COLOR_DARK = RGBColor(20, 20, 20)
    COLOR_MUTED = RGBColor(80, 80, 80)
    
    # 1. Iterate through all slides and remove all empty placeholders and ghost textboxes
    for s_idx, slide in enumerate(prs.slides):
        shapes_to_delete = []
        for shape in slide.shapes:
            txt = shape.text.strip() if hasattr(shape, 'text') else ''
            
            # Delete empty placeholders (these cause "Click to add title", "Moment/Topic", "Reinvention message")
            if shape.is_placeholder and not txt:
                shapes_to_delete.append(shape)
            # Delete non-placeholder textboxes that are completely empty
            elif hasattr(shape, 'text') and not txt and shape.shape_type == 17: # 17 is TEXT_BOX
                # Only delete if it has no fill or border
                if not shape.fill.type or shape.fill.type == 0:
                    shapes_to_delete.append(shape)
                    
        for shape in shapes_to_delete:
            sp = shape._element
            sp.getparent().remove(sp)
            
    # 2. Fix top slide titles and purple banner formatting on each content slide
    # On Slides 3, 4, 5, 6, 7, 8:
    for s_idx in range(2, len(prs.slides) - 1):
        slide = prs.slides[s_idx]
        
        # Find the purple banner shape and the title textbox
        purple_bar = None
        title_box = None
        
        for shape in slide.shapes:
            if shape.name.startswith("Rectangle 1") or (shape.top < Inches(1.0) and shape.height > Inches(0.4) and shape.width > Inches(8.0) and shape.shape_type == 1):
                purple_bar = shape
            elif hasattr(shape, 'text') and shape.top < Inches(1.2) and shape.width > Inches(6.0):
                txt = shape.text.strip()
                if any(txt.startswith(p) for p in ["The problem:", "Foundation:", "The 3-model", "Designed for", "From root cause", "Enterprise-ready"]):
                    title_box = shape

        # If purple bar exists, position title box cleanly inside or style it
        if purple_bar and title_box:
            # Set purple bar exact coordinates
            purple_bar.left = Inches(0.5)
            purple_bar.top = Inches(0.3)
            purple_bar.width = Inches(12.333)
            purple_bar.height = Inches(0.65)
            
            # Position title box inside the purple bar with white text
            title_box.left = Inches(0.7)
            title_box.top = Inches(0.32)
            title_box.width = Inches(11.9)
            title_box.height = Inches(0.6)
            
            tf = title_box.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.0)
            tf.margin_top = Inches(0.05)
            for p in tf.paragraphs:
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = COLOR_WHITE
        elif title_box:
            # If no purple bar, give title clean top position
            title_box.top = Inches(0.4)
            title_box.left = Inches(0.6)
            tf = title_box.text_frame
            for p in tf.paragraphs:
                p.font.size = Pt(22)
                p.font.bold = True
                p.font.color.rgb = COLOR_PURPLE

    # 3. Slide 6 (Scenario 4) formatting fix:
    # Ensure Scenario 4 title and description are clean and not overlapping
    s6 = prs.slides[5]
    for shape in s6.shapes:
        if hasattr(shape, 'text'):
            txt = shape.text.strip()
            if txt.startswith("SCENARIO 4"):
                shape.top = Inches(4.7)
                shape.left = Inches(0.6)
                shape.width = Inches(12.0)
                for p in shape.text_frame.paragraphs:
                    p.font.size = Pt(14)
                    p.font.bold = True
                    p.font.color.rgb = COLOR_PURPLE
            elif txt.startswith("Sensitive cost / margin"):
                shape.top = Inches(5.1)
                shape.left = Inches(0.6)
                shape.width = Inches(12.0)
                shape.height = Inches(1.5)
                for p in shape.text_frame.paragraphs:
                    p.font.size = Pt(13)
                    p.font.color.rgb = COLOR_DARK

    # 4. Slide 7 (Recovery Plan) formatting fix:
    s7 = prs.slides[6]
    for shape in s7.shapes:
        if hasattr(shape, 'text'):
            txt = shape.text.strip()
            if txt.startswith("PERSONA-TAILORED"):
                shape.top = Inches(3.9)
                shape.left = Inches(0.6)
                for p in shape.text_frame.paragraphs:
                    p.font.size = Pt(14)
                    p.font.bold = True
                    p.font.color.rgb = COLOR_PURPLE
            elif txt.startswith("Human-in-the-loop feedback"):
                shape.top = Inches(6.3)
                shape.left = Inches(0.6)
                shape.width = Inches(12.0)
                for p in shape.text_frame.paragraphs:
                    p.font.size = Pt(12)
                    p.font.color.rgb = COLOR_MUTED

    # 5. Slide 8 (Enterprise Controls) formatting fix:
    s8 = prs.slides[7]
    for shape in s8.shapes:
        if hasattr(shape, 'text'):
            txt = shape.text.strip()
            if txt.startswith("Core principle:"):
                shape.top = Inches(6.3)
                shape.left = Inches(0.6)
                shape.width = Inches(12.0)
                for p in shape.text_frame.paragraphs:
                    p.font.size = Pt(13)
                    p.font.bold = True
                    p.font.color.rgb = COLOR_PURPLE

    # Save outputs
    for out_fn in [output_file, "Warriors_BusinessIntelligence_Accenture_2026_FINAL.pptx", "Warriors_BusinessIntelligence_Accenture_2026_FONT_FIXED.pptx"]:
        try:
            prs.save(out_fn)
            print(f"Cleaned & perfected deck saved to: {out_fn}")
        except PermissionError:
            print(f"Note: {out_fn} is currently open in PowerPoint on desktop.")

if __name__ == "__main__":
    in_fn = "Warriors_BusinessIntelligence_Accenture_2026_FONT_FIXED.pptx"
    out_fn = "Warriors_BusinessIntelligence_Accenture_2026_CLEAN.pptx"
    fix_presentation(in_fn, out_fn)
