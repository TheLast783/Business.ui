"""
Upgrade and redesign Warriors_ 3-Model Triad.pptx using the official Accenture template.
Preserves Slides 1 & 2 intact, and upgrades Slides 3-6 with high-legibility cards, medium content, and bigger fonts.
"""
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def redesign_warriors_deck(input_path: str, output_path: str):
    prs = Presentation(input_path)
    
    # Palette definition for cards (designed to fit over the template)
    COLOR_CARD = RGBColor(15, 23, 42)         # Slate Navy #0F172A
    COLOR_BORDER = RGBColor(30, 41, 59)       # Border #1E293B
    COLOR_PRIMARY = RGBColor(2, 132, 199)     # Sky Blue #0284C7
    COLOR_ACCENT = RGBColor(56, 189, 248)     # Accent Blue #38BDF8
    COLOR_TEXT = RGBColor(248, 250, 252)      # Bright White #F8FAFC
    COLOR_MUTED = RGBColor(148, 163, 184)     # Slate Gray #94A3B8
    COLOR_GREEN = RGBColor(16, 185, 129)      # Emerald #10B981
    COLOR_AMBER = RGBColor(245, 158, 11)      # Amber #F59E0B

    def clear_content_shapes(slide):
        """Remove existing body textboxes while keeping slide title / background elements."""
        shapes_to_remove = []
        for sh in slide.shapes:
            # Keep shapes that look like slide titles or header background
            if hasattr(sh, 'text') and sh.text.strip():
                txt = sh.text.strip()
                if txt.startswith("Describe the problem") or txt.startswith("Proposed solution") or txt.startswith("Thank you"):
                    continue
                shapes_to_remove.append(sh)
            elif sh.name.startswith("TextBox") or sh.name.startswith("Text Placeholder"):
                shapes_to_remove.append(sh)
        for sh in shapes_to_remove:
            sp = sh._element
            sp.getparent().remove(sp)

    def add_card(slide, left, top, width, height, bg_color=COLOR_CARD, border_color=COLOR_BORDER):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
        return shape

    # -------------------------------------------------------------
    # SLIDE 3: PROBLEM STATEMENT (200 words)
    # -------------------------------------------------------------
    s3 = prs.slides[2]
    clear_content_shapes(s3)
    
    # Update title text
    for sh in s3.shapes:
        if hasattr(sh, 'text') and "problem" in sh.text.lower():
            sh.text = "Describe the problem statement (200 words)"
            for p in sh.text_frame.paragraphs:
                p.font.size = Pt(22)
                p.font.bold = True
                p.font.color.rgb = COLOR_TEXT

    # Left Card: Core Challenge
    add_card(s3, Inches(0.6), Inches(1.3), Inches(5.8), Inches(5.5))
    tb3_1 = s3.shapes.add_textbox(Inches(0.85), Inches(1.5), Inches(5.3), Inches(5.1))
    tf3_1 = tb3_1.text_frame
    tf3_1.word_wrap = True

    p = tf3_1.paragraphs[0]
    p.text = "The Operational & Analytical Dilemma"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_AMBER

    s3_left_points = [
        ("Static Dashboards Fail at 'Why': ", "Modern dashboards report what happened with red alerts, but fail to explain why or what next. Leaders cannot differentiate internal warehouse glitches from macroeconomic market shocks."),
        ("Manual Multi-Day Triage: ", "Root cause isolation falls on human analysts spending 4 to 7 days manually writing SQL queries across ERP sales tables, web event logs, and support tickets.")
    ]
    for b_title, b_desc in s3_left_points:
        p = tf3_1.add_paragraph()
        p.space_before = Pt(14)
        r1 = p.add_run()
        r1.text = b_title
        r1.font.bold = True
        r1.font.size = Pt(14)
        r1.font.color.rgb = COLOR_TEXT
        r2 = p.add_run()
        r2.text = b_desc
        r2.font.size = Pt(13)
        r2.font.color.rgb = COLOR_MUTED

    # Right Card: The 3 Critical Risks
    add_card(s3, Inches(6.8), Inches(1.3), Inches(5.9), Inches(5.5))
    tb3_2 = s3.shapes.add_textbox(Inches(7.05), Inches(1.5), Inches(5.4), Inches(5.1))
    tf3_2 = tb3_2.text_frame
    tf3_2.word_wrap = True

    p = tf3_2.paragraphs[0]
    p.text = "Three Critical Enterprise Bottlenecks"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    s3_right_points = [
        ("1. Critical Latency: ", "By the time manual root cause is diagnosed days later, customer churn and revenue loss become permanent."),
        ("2. Severe Alert Fatigue: ", "Dashboards trigger false alarms on normal cyclical/weekend swings because they lack statistical process control."),
        ("3. AI Hallucination & Data Leakage: ", "Generic LLMs fabricate arithmetic calculations and risk exposing sensitive cost margins across internal roles.")
    ]
    for b_title, b_desc in s3_right_points:
        p = tf3_2.add_paragraph()
        p.space_before = Pt(14)
        r1 = p.add_run()
        r1.text = b_title
        r1.font.bold = True
        r1.font.size = Pt(14)
        r1.font.color.rgb = COLOR_TEXT
        r2 = p.add_run()
        r2.text = b_desc
        r2.font.size = Pt(13)
        r2.font.color.rgb = COLOR_MUTED

    # -------------------------------------------------------------
    # SLIDE 4: PROPOSED SOLUTION (PART 1 - MATH & DIAGNOSTICS)
    # -------------------------------------------------------------
    s4 = prs.slides[3]
    clear_content_shapes(s4)

    for sh in s4.shapes:
        if hasattr(sh, 'text') and "solution" in sh.text.lower():
            sh.text = "Proposed solution (200 words) — Diagnostic Engine"
            for p in sh.text_frame.paragraphs:
                p.font.size = Pt(22)
                p.font.bold = True
                p.font.color.rgb = COLOR_TEXT

    # Left: Deterministic Math Core
    add_card(s4, Inches(0.6), Inches(1.3), Inches(5.8), Inches(5.5), border_color=COLOR_PRIMARY)
    tb4_1 = s4.shapes.add_textbox(Inches(0.85), Inches(1.5), Inches(5.3), Inches(5.1))
    tf4_1 = tb4_1.text_frame
    tf4_1.word_wrap = True

    p = tf4_1.paragraphs[0]
    p.text = "1. Deterministic Non-LLM Math Core"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    s4_left_points = [
        ("Statistical Process Control (SPC): ", "Applies 28-day Day-of-Week normalized baseline with a >2.5sigma threshold to filter normal seasonal variance from true anomalies with 0 LLM tokens."),
        ("Exact 3-Factor Shapley Metric Tree: ", "Mathematically decomposes revenue changes into Volume, Conversion Rate, and AOV with strict $0.00 residual error (|eps| < 10^-12).")
    ]
    for b_title, b_desc in s4_left_points:
        p = tf4_1.add_paragraph()
        p.space_before = Pt(14)
        r1 = p.add_run()
        r1.text = b_title
        r1.font.bold = True
        r1.font.size = Pt(14)
        r1.font.color.rgb = COLOR_TEXT
        r2 = p.add_run()
        r2.text = b_desc
        r2.font.size = Pt(13)
        r2.font.color.rgb = COLOR_MUTED

    # Right: Dual Diagnostic Intelligences
    add_card(s4, Inches(6.8), Inches(1.3), Inches(5.9), Inches(5.5), border_color=COLOR_GREEN)
    tb4_2 = s4.shapes.add_textbox(Inches(7.05), Inches(1.5), Inches(5.4), Inches(5.1))
    tf4_2 = tb4_2.text_frame
    tf4_2.word_wrap = True

    p = tf4_2.paragraphs[0]
    p.text = "2. Dual Diagnostic Intelligences"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN

    s4_right_points = [
        ("Model 1 (Internal Operational Specialist): ", "Operates on-premise over ERP staging queues and Jira tickets (JIRA-4819) to isolate internal backlogs with zero external data exposure."),
        ("Model 2 (Live Global Sentinel): ", "Ingests real-time external maritime feeds (MACRO-PORT-01) to quantify port congestion and external macro shocks invisible to internal logs.")
    ]
    for b_title, b_desc in s4_right_points:
        p = tf4_2.add_paragraph()
        p.space_before = Pt(14)
        r1 = p.add_run()
        r1.text = b_title
        r1.font.bold = True
        r1.font.size = Pt(14)
        r1.font.color.rgb = COLOR_TEXT
        r2 = p.add_run()
        r2.text = b_desc
        r2.font.size = Pt(13)
        r2.font.color.rgb = COLOR_MUTED

    # -------------------------------------------------------------
    # SLIDE 5: PROPOSED SOLUTION (PART 2 - SIMULATION & GOVERNANCE)
    # -------------------------------------------------------------
    s5 = prs.slides[4]
    clear_content_shapes(s5)

    for sh in s5.shapes:
        if hasattr(sh, 'text') and "solution" in sh.text.lower():
            sh.text = "Proposed solution (200 words) — Simulation & Action"
            for p in sh.text_frame.paragraphs:
                p.font.size = Pt(22)
                p.font.bold = True
                p.font.color.rgb = COLOR_TEXT

    # Left: Model 3 & Simulation
    add_card(s5, Inches(0.6), Inches(1.3), Inches(5.8), Inches(5.5), border_color=COLOR_GREEN)
    tb5_1 = s5.shapes.add_textbox(Inches(0.85), Inches(1.5), Inches(5.3), Inches(5.1))
    tf5_1 = tb5_1.text_frame
    tf5_1.word_wrap = True

    p = tf5_1.paragraphs[0]
    p.text = "3. Model 3: Prescriptive & ROI Simulator"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN

    s5_left_points = [
        ("Attribution Synthesis: ", "Reconciles Model 1 and Model 2 findings into exact causal attribution shares (e.g. 70% port strike + 30% warehouse backlog)."),
        ("Pre-Action Outcome Simulation: ", "Forecasts 30/60/90-day recovery curves across Status Quo vs Prescribed Recovery (4.2x projected ROI) before spending capital."),
        ("Dynamic Executive Controls: ", "Sliders for remediation budget cap ($45k) and policy overrides instantly recalculate trajectory outcomes in real time.")
    ]
    for b_title, b_desc in s5_left_points:
        p = tf5_1.add_paragraph()
        p.space_before = Pt(12)
        r1 = p.add_run()
        r1.text = b_title
        r1.font.bold = True
        r1.font.size = Pt(14)
        r1.font.color.rgb = COLOR_TEXT
        r2 = p.add_run()
        r2.text = b_desc
        r2.font.size = Pt(13)
        r2.font.color.rgb = COLOR_MUTED

    # Right: Safety, Abstention & RBAC
    add_card(s5, Inches(6.8), Inches(1.3), Inches(5.9), Inches(5.5), border_color=COLOR_AMBER)
    tb5_2 = s5.shapes.add_textbox(Inches(7.05), Inches(1.5), Inches(5.4), Inches(5.1))
    tf5_2 = tb5_2.text_frame
    tf5_2.word_wrap = True

    p = tf5_2.paragraphs[0]
    p.text = "4. AI Safety & Governed Security"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_AMBER

    s5_right_points = [
        ("Explicit Abstention Protocol: ", "When evidence confidence delta is <25%, the engine halts high-capital interventions and prescribes low-cost canary validation tests ($150 / 2 hours)."),
        ("Persona Tailoring & RBAC: ", "Delivers strategic briefs for Executives while generating tactical SOPs and redacting confidential unit COGS and margins for Operations Analysts.")
    ]
    for b_title, b_desc in s5_right_points:
        p = tf5_2.add_paragraph()
        p.space_before = Pt(14)
        r1 = p.add_run()
        r1.text = b_title
        r1.font.bold = True
        r1.font.size = Pt(14)
        r1.font.color.rgb = COLOR_TEXT
        r2 = p.add_run()
        r2.text = b_desc
        r2.font.size = Pt(13)
        r2.font.color.rgb = COLOR_MUTED

    # -------------------------------------------------------------
    # SLIDE 6: THANK YOU & VERIFICATION HIGHLIGHTS
    # -------------------------------------------------------------
    s6 = prs.slides[5]
    clear_content_shapes(s6)

    for sh in s6.shapes:
        if hasattr(sh, 'text') and "thank" in sh.text.lower():
            sh.text = "Delivered Working Prototype & Verification Highlights"
            for p in sh.text_frame.paragraphs:
                p.font.size = Pt(22)
                p.font.bold = True
                p.font.color.rgb = COLOR_TEXT

    # 4 Cards Matrix on Slide 6
    cards_s6 = [
        ("Interactive Decision Workspace", "Live Python + Streamlit application with scenario selector, persona switcher, Plotly SPC control charts, and 30/60/90-day trajectory curves.", Inches(0.6), Inches(1.3)),
        ("93 / 93 Automated Passing Tests", "Comprehensive test suite verifying zero residual error (|eps| < 10^-12), zero LLM math tokens, and end-to-end scenario execution.", Inches(6.8), Inches(1.3)),
        ("All 4 Required Scenarios Covered", "1. Multi-factor compound drop (70/30)\n2. Low-confidence ambiguity with canary tests\n3. Sparse-history cold start launch\n4. Role-based data masking.", Inches(0.6), Inches(4.0)),
        ("Pluggable & Secure Architecture", "Instant toggle between live API providers (Gemini, OpenAI, Ollama) and offline deterministic fallback (<5ms latency, $0.00 cost).", Inches(6.8), Inches(4.0))
    ]

    for title, desc, left_pos, top_pos in cards_s6:
        add_card(s6, left_pos, top_pos, Inches(5.9), Inches(2.3))
        tb = s6.shapes.add_textbox(left_pos + Inches(0.2), top_pos + Inches(0.15), Inches(5.5), Inches(2.0))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT

        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = COLOR_MUTED
        p_desc.space_before = Pt(6)

    # Footer signature
    tb_sig = s6.shapes.add_textbox(Inches(0.6), Inches(6.5), Inches(12.0), Inches(0.6))
    tf_sig = tb_sig.text_frame
    p_sig = tf_sig.paragraphs[0]
    p_sig.text = "Team: Warriors (IIT Patna) — Harsh Singh Baghel, Rathod Rudra, Dhruv Maheshwari"
    p_sig.font.size = Pt(13)
    p_sig.font.bold = True
    p_sig.font.color.rgb = COLOR_TEXT

    prs.save(output_path)
    print(f"Redesigned presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    in_file = "Warriors_ 3-Model Triad.pptx"
    out_file = "Warriors_ 3-Model Triad.pptx"
    redesign_warriors_deck(in_file, out_file)
