"""
Generate Detailed Business Proposal in PPTX and PDF for the Accenture Innovation Challenge 2026 Submission.
Formats Slides 3 to 6 strictly as README & Implementation Architecture documentation.
"""
import os
import sys
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def build_presentation():
    input_pptx = "Warriors_ 3-Model Triad_backup.pptx"
    output_pptx = "Detailed_Business_Proposal_Warriors.pptx"
    output_pdf = "Detailed_Business_Proposal_Warriors.pdf"

    prs = pptx.Presentation(input_pptx)
    print(f"Loaded {input_pptx} with {len(prs.slides)} slides.")

    purple_title = RGBColor(70, 20, 130)      # Accenture Deep Purple
    purple_accent = RGBColor(161, 0, 255)     # Accenture Purple Accent
    blue_accent = RGBColor(0, 102, 204)       # Tech Blue
    dark_body = RGBColor(40, 40, 40)          # Charcoal Text
    sub_body = RGBColor(60, 60, 60)           # Secondary Text

    # =========================================================================
    # SLIDE 3: README — Implementation Approach & Methodology
    # =========================================================================
    slide3 = prs.slides[2]
    for shape in slide3.shapes:
        if "Title" in shape.name or (shape.has_text_frame and "problem statement" in shape.text_frame.text.lower()):
            shape.text_frame.text = "README: Implementation Approach & Methodology"
            if shape.text_frame.paragraphs:
                p = shape.text_frame.paragraphs[0]
                p.font.name = "Arial"
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = purple_title
        elif shape.has_text_frame and shape.name != "Title 3":
            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True

            p0 = tf.paragraphs[0]
            p0.text = "1. Multi-Source Ingestion & Governed Canonical Contract:"
            p0.font.name = "Arial"
            p0.font.size = Pt(10.5)
            p0.font.bold = True
            p0.font.color.rgb = purple_title
            p0.space_after = Pt(2)

            p0_desc = tf.add_paragraph()
            p0_desc.text = "Harmonizes 3 mismatched enterprise grains (Daily ERP sales SQL, Hourly Web clickstream sessions, and Weekly Jira support tickets) into a unified daily canonical snapshot via MultiSourceDataLoader."
            p0_desc.font.name = "Arial"
            p0_desc.font.size = Pt(9.5)
            p0_desc.font.color.rgb = sub_body
            p0_desc.space_after = Pt(5)

            p1 = tf.add_paragraph()
            p1.text = "2. Deterministic Non-LLM Mathematical Core (Zero Hallucinations):"
            p1.font.name = "Arial"
            p1.font.size = Pt(10.5)
            p1.font.bold = True
            p1.font.color.rgb = purple_title
            p1.space_after = Pt(2)

            bullets1 = [
                ("SPC Deseasonalization: ", "28-day rolling Day-of-Week (DoW) normalized baseline with 2.5-sigma threshold mathematically filters cyclical noise before alerting."),
                ("Exact Shapley Attribution: ", "Closed-form 3-factor metric tree decomposition (Revenue = Sessions * CVR * AOV) with guaranteed 0.00 residual error (|eps| < 10^-12) consuming exactly 0 LLM math tokens.")
            ]
            for title, desc in bullets1:
                pr = tf.add_paragraph()
                pr.space_after = Pt(2)
                run1 = pr.add_run()
                run1.text = "- " + title
                run1.font.name = "Arial"
                run1.font.size = Pt(9.5)
                run1.font.bold = True
                run1.font.color.rgb = blue_accent
                run2 = pr.add_run()
                run2.text = desc
                run2.font.name = "Arial"
                run2.font.size = Pt(9.5)
                run2.font.color.rgb = sub_body

            p2 = tf.add_paragraph()
            p2.text = "3. Triangulated AI Synthesis, Abstention & Security:"
            p2.font.name = "Arial"
            p2.font.size = Pt(10.5)
            p2.font.bold = True
            p2.font.color.rgb = purple_title
            p2.space_before = Pt(3)
            p2.space_after = Pt(2)

            bullets2 = [
                ("3-Tier AI Loop: ", "Model 1 diagnoses on-premise Jira/ERP logs; Model 2 monitors live macro feeds; Model 3 performs attribution % and 30/60/90-day trajectory ROI simulation."),
                ("Explicit Abstention Protocol: ", "When competing cause margin is <25%, halts high-capital action and prescribes 2 low-cost canary tests ($120-$350)."),
                ("Role-Based Access Control (RBAC): ", "Masks sensitive unit COGS and gross margins for Operations Analysts while granting full access to Executives.")
            ]
            for title, desc in bullets2:
                pr = tf.add_paragraph()
                pr.space_after = Pt(2)
                run1 = pr.add_run()
                run1.text = "- " + title
                run1.font.name = "Arial"
                run1.font.size = Pt(9.5)
                run1.font.bold = True
                run1.font.color.rgb = purple_accent
                run2 = pr.add_run()
                run2.text = desc
                run2.font.name = "Arial"
                run2.font.size = Pt(9.5)
                run2.font.color.rgb = sub_body

    # =========================================================================
    # SLIDE 4: README — Solution Architecture & Engine Layers
    # =========================================================================
    slide4 = prs.slides[3]
    for shape in slide4.shapes:
        if "Title" in shape.name or (shape.has_text_frame and "proposed solution" in shape.text_frame.text.lower()):
            shape.text_frame.text = "README: Solution Architecture & Engine Layers"
            if shape.text_frame.paragraphs:
                p = shape.text_frame.paragraphs[0]
                p.font.name = "Arial"
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = purple_title
        elif shape.has_text_frame and shape.name != "Title 3":
            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True

            arch_layers = [
                ("Layer 1: Governed Ingestion & Schema Contracts (engine/contracts/, engine/data/)", [
                    ("Canonical Data Harmonization: ", "Ingests daily ERP sales, hourly web sessions, and weekly Jira tickets into validated Pydantic schemas."),
                    ("Dynamic RBAC Masking Engine: ", "Performs real-time column/row redaction on financial margins based on active user persona token.")
                ]),
                ("Layer 2: Deterministic Non-LLM Math Core (engine/math/)", [
                    ("Statistical Process Control (spc.py): ", "28-day DoW baseline, Student-t cold-start envelope (N<14), and MAD robust outlier detection."),
                    ("Axiomatic Causal Metric Tree (causal_tree.py): ", "Exact 3-factor Shapley attribution (Sessions, CVR, AOV) + LMDI-1 validation.")
                ]),
                ("Layer 3: 3-Model AI Synthesis & Simulation (engine/synthesis/)", [
                    ("Model 1 (Internal Diagnostic): ", "Private on-premise LLM diagnosing warehouse backlogs and Jira incident tickets."),
                    ("Model 2 (Macro Sentinel): ", "Live API sentinel ingesting port strikes, freight indices, and competitor campaigns."),
                    ("Model 3 (Prescriptive Trajectory Simulator): ", "Attribution weighting + 30/60/90-day ROI curves under executive budget caps ($45k)."),
                    ("Abstention Engine (abstention.py): ", "Calculates hypothesis confidence delta; triggers canary validation tests if margin < 25%.")
                ]),
                ("Layer 4: Interactive Decision Workspace & Telemetry (ui/, engine/telemetry/)", [
                    ("Streamlit UI & Telemetry Tracker: ", "6-tab workspace with Plotly charts, persona switchers, human feedback loop, and ms latency/token accounting.")
                ])
            ]

            is_first = True
            for header, items in arch_layers:
                p_hdr = tf.paragraphs[0] if is_first else tf.add_paragraph()
                is_first = False
                p_hdr.text = header
                p_hdr.font.name = "Arial"
                p_hdr.font.size = Pt(10)
                p_hdr.font.bold = True
                p_hdr.font.color.rgb = purple_title
                p_hdr.space_before = Pt(2)
                p_hdr.space_after = Pt(1)

                for b_title, b_desc in items:
                    pr = tf.add_paragraph()
                    pr.space_after = Pt(1.5)
                    run1 = pr.add_run()
                    run1.text = "- " + b_title
                    run1.font.name = "Arial"
                    run1.font.size = Pt(9)
                    run1.font.bold = True
                    run1.font.color.rgb = blue_accent
                    run2 = pr.add_run()
                    run2.text = b_desc
                    run2.font.name = "Arial"
                    run2.font.size = Pt(9)
                    run2.font.color.rgb = sub_body

    # =========================================================================
    # SLIDE 5: README — Dependencies & Technology Stack
    # =========================================================================
    slide5 = prs.slides[4]
    for shape in slide5.shapes:
        if "Title" in shape.name or (shape.has_text_frame and "proposed solution" in shape.text_frame.text.lower()):
            shape.text_frame.text = "README: Dependencies & Technology Ecosystem"
            if shape.text_frame.paragraphs:
                p = shape.text_frame.paragraphs[0]
                p.font.name = "Arial"
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = purple_title
        elif shape.has_text_frame and "Title" not in shape.name:
            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True

            dep_sections = [
                ("Runtime Core & Execution Environment:", [
                    ("Python 3.11+ / PyPy: ", "Modern Python runtime utilizing typing, dataclasses, and async execution."),
                    ("Pydantic v2 (>=2.5.0): ", "Strict runtime schema enforcement, data contract validation, and JSON serialization.")
                ]),
                ("Front-End Interactive Workspace & Visualization:", [
                    ("Streamlit (>=1.30.0): ", "Interactive multi-tab executive dashboard framework with session state management."),
                    ("Plotly (>=5.18.0): ", "Interactive WebGL/SVG visualization engine powering SPC Control Charts, Shapley Waterfall, and 30/60/90-Day Trajectory ROI curves.")
                ]),
                ("Data Processing & Mathematical Computation:", [
                    ("Pandas (>=2.1.0): ", "High-performance time-series alignment, daily grain resampling, and DataFrame RBAC masking."),
                    ("NumPy (>=1.26.0): ", "Vectorized numerical routines, Student-t distribution limits, and Median Absolute Deviation (MAD).")
                ]),
                ("Pluggable Multi-LLM Provider Architecture:", [
                    ("Hybrid Provider Interface: ", "Pluggable provider adapter supporting Google Gemini, OpenAI GPT-4o, and local Ollama."),
                    ("Deterministic Offline Mock Fallback: ", "Built-in deterministic mock provider for instant execution with <5ms latency and $0.00 token cost.")
                ]),
                ("Verification & Testing Infrastructure:", [
                    ("Python unittest & Monte Carlo: ", "Headless test harness executing 142 assertions and 10,000 randomized invariant stress trials.")
                ])
            ]

            is_first = True
            for header, items in dep_sections:
                p_hdr = tf.paragraphs[0] if is_first else tf.add_paragraph()
                is_first = False
                p_hdr.text = header
                p_hdr.font.name = "Arial"
                p_hdr.font.size = Pt(10)
                p_hdr.font.bold = True
                p_hdr.font.color.rgb = purple_title
                p_hdr.space_before = Pt(2)
                p_hdr.space_after = Pt(1)

                for b_title, b_desc in items:
                    pr = tf.add_paragraph()
                    pr.space_after = Pt(1.5)
                    run1 = pr.add_run()
                    run1.text = "- " + b_title
                    run1.font.name = "Arial"
                    run1.font.size = Pt(9)
                    run1.font.bold = True
                    run1.font.color.rgb = purple_accent
                    run2 = pr.add_run()
                    run2.text = b_desc
                    run2.font.name = "Arial"
                    run2.font.size = Pt(9)
                    run2.font.color.rgb = sub_body

    # =========================================================================
    # SLIDE 6: README — Execution Instructions & Test Verification
    # =========================================================================
    slide6 = prs.slides[5]
    for shape in slide6.shapes:
        if "Title" in shape.name or (shape.has_text_frame and "thank you" in shape.text_frame.text.lower()):
            shape.text_frame.text = "README: Execution Instructions & Test Verification"
            if shape.text_frame.paragraphs:
                p = shape.text_frame.paragraphs[0]
                p.font.name = "Arial"
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = purple_title

    slide6_textbox = None
    for shape in slide6.shapes:
        if shape.has_text_frame and "Title" not in shape.name:
            slide6_textbox = shape
            break
    
    if not slide6_textbox:
        left = Inches(1.0)
        top = Inches(2.0)
        width = Inches(11.3)
        height = Inches(4.8)
        slide6_textbox = slide6.shapes.add_textbox(left, top, width, height)

    tf6 = slide6_textbox.text_frame
    tf6.clear()
    tf6.word_wrap = True

    exec_sections = [
        ("1. Environment Setup & Dependency Installation:", [
            ("Terminal Command: ", "pip install -r prototype/requirements.txt"),
            ("Prerequisites: ", "Python 3.11+ installed; virtual environment recommended.")
        ]),
        ("2. Launch Interactive Decision Workspace UI:", [
            ("Terminal Command: ", "streamlit run prototype/app.py"),
            ("Browser Access: ", "Open http://localhost:8501 (Interactive scenarios, persona toggles, trajectory sliders).")
        ]),
        ("3. Run Automated Headless Verification Test Suite:", [
            ("Acceptance Tests (93/93 Passed): ", "python prototype/test_scenarios.py"),
            ("Modular Unit & Property Tests (49/49 Passed): ", "python -m unittest discover -s prototype/tests -p \"test_*.py\""),
            ("Scenario x Persona Matrix Runner (8/8 Passed): ", "python prototype/verify_e2e.py")
        ]),
        ("4. Key Verification Guarantees & Acceptance Matrix:", [
            ("Deterministic Math Invariant: ", "Shapley metric tree residual error 0.00000000 (|eps| < 10^-12) across 10,000 Monte Carlo trials."),
            ("Zero Math Token Invariant: ", "Statistical Process Control and Causal Tree attribution consume exactly 0 LLM tokens ($0.00 cost)."),
            ("All 4 Mandatory Scenarios: ", "1. Multi-factor (70/30) | 2. Ambiguity & Canary tests | 3. Cold start prior | 4. RBAC masking.")
        ])
    ]

    is_first = True
    for header, items in exec_sections:
        p_hdr = tf6.paragraphs[0] if is_first else tf6.add_paragraph()
        is_first = False
        p_hdr.text = header
        p_hdr.font.name = "Arial"
        p_hdr.font.size = Pt(10)
        p_hdr.font.bold = True
        p_hdr.font.color.rgb = purple_title
        p_hdr.space_before = Pt(2)
        p_hdr.space_after = Pt(1)

        for b_title, b_desc in items:
            pr = tf6.add_paragraph()
            pr.space_after = Pt(1.5)
            run1 = pr.add_run()
            run1.text = "- " + b_title
            run1.font.name = "Arial"
            run1.font.size = Pt(9)
            run1.font.bold = True
            run1.font.color.rgb = blue_accent
            run2 = pr.add_run()
            run2.text = b_desc
            run2.font.name = "Arial"
            run2.font.size = Pt(9)
            run2.font.color.rgb = sub_body

    # Save PPTX output
    prs.save(output_pptx)
    print(f"Saved: {output_pptx}")

    # Also try to overwrite original if not locked
    try:
        prs.save("Warriors_ 3-Model Triad.pptx")
        print("Updated: Warriors_ 3-Model Triad.pptx")
    except Exception as e:
        print(f"Note: Warriors_ 3-Model Triad.pptx is open in PowerPoint. Saved to {output_pptx}")

    # Convert to PDF via win32com PowerPoint COM automation
    try:
        import win32com.client
        import pythoncom
        pythoncom.CoInitialize()
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        abs_pptx = os.path.abspath(output_pptx)
        abs_pdf = os.path.abspath(output_pdf)
        
        # Open and export
        deck = powerpoint.Presentations.Open(abs_pptx, WithWindow=False)
        # Format 32 = ppSaveAsPDF
        deck.SaveAs(abs_pdf, 32)
        deck.Close()
        powerpoint.Quit()
        print(f"Successfully converted to PDF: {output_pdf}")
    except Exception as e:
        print(f"PowerPoint COM export note: {e}")

if __name__ == "__main__":
    build_presentation()
