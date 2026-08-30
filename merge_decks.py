"""
Merge the original template (Slide 1, Slide 2, and OG Last Slide) from Warriors_ 3-Model Triad_backup.pptx
with the 9 comprehensive, high-legibility cards slides into Warriors_ 3-Model Triad.pptx.
"""
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_merged_warriors_deck(backup_path: str, output_path: str):
    # 1. Load the original presentation template
    prs = Presentation(backup_path)
    
    # Extract OG Slide 1, OG Slide 2, and OG Slide 6 (the last slide)
    # In python-pptx, we can delete slides 2, 3, 4 (0-indexed: 2, 3, 4) or build fresh slides using the blank/content layout
    
    # Color definitions
    COLOR_BG = RGBColor(9, 14, 23)          # #090E17
    COLOR_CARD = RGBColor(15, 23, 42)       # #0F172A
    COLOR_BORDER = RGBColor(30, 41, 59)     # #1E293B
    COLOR_PRIMARY = RGBColor(2, 132, 199)   # #0284C7
    COLOR_ACCENT = RGBColor(56, 189, 248)   # #38BDF8
    COLOR_TEXT = RGBColor(248, 250, 252)    # #F8FAFC
    COLOR_MUTED = RGBColor(148, 163, 184)   # #94A3B8
    COLOR_GREEN = RGBColor(16, 185, 129)    # #10B981
    COLOR_AMBER = RGBColor(245, 158, 11)    # #F59E0B

    # Helper to create styled cards
    def add_card(slide, left, top, width, height, bg_color=COLOR_CARD, border_color=COLOR_BORDER):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
        return shape

    def add_header(slide, title_text: str, category_text: str = "BUSINESSINTELLIGENCE.AI · ROUND 2"):
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.35))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_ACCENT

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.7), Inches(0.75))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT

    def set_dark_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.color.rgb = COLOR_BG
        return bg

    # Let's remove the old middle slides (slides at index 2, 3, 4)
    # Delete from back to front: index 4, then index 3, then index 2
    rId4 = prs.slides._sldIdLst[4].rId
    prs.part.drop_rel(rId4)
    del prs.slides._sldIdLst[4]

    rId3 = prs.slides._sldIdLst[3].rId
    prs.part.drop_rel(rId3)
    del prs.slides._sldIdLst[3]

    rId2 = prs.slides._sldIdLst[2].rId
    prs.part.drop_rel(rId2)
    del prs.slides._sldIdLst[2]

    # Now prs has 3 slides: Slide 1 (Cover), Slide 2 (Team Details), Slide 3 (OG Thank You)
    # We will insert the 9 slides before the last slide (or create slides and move the last slide)
    blank_layout = prs.slide_layouts[6]

    # Function to build slide content
    def build_middle_slides():
        # --- Slide 1/9: Executive Problem Statement ---
        s = prs.slides.add_slide(blank_layout)
        set_dark_bg(s)
        add_header(s, "Executive Problem: The Telemetry-to-Action Gap")
        add_card(s, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.2))
        tb1 = s.shapes.add_textbox(Inches(1.05), Inches(1.9), Inches(5.1), Inches(4.7))
        tf1 = tb1.text_frame
        tf1.word_wrap = True
        p = tf1.paragraphs[0]
        p.text = "The Operational Dilemma"
        p.font.size = Pt(19)
        p.font.bold = True
        p.font.color.rgb = COLOR_AMBER
        for b_t, n_t in [
            ("Noisy Anomaly Alerts: ", "Static dashboards flood leaders with false alarms, failing to isolate seasonal Day-of-Week swings from genuine operational breaks."),
            ("Multi-Source Data Silos: ", "ERP SQL ledgers (daily), web session telemetry (hourly), and Jira tickets (weekly) are disconnected."),
            ("The LLM Hallucination Risk: ", "Pure LLM analytics fabricate arithmetic totals, violate financial identities, and propose ungrounded actions.")
        ]:
            p = tf1.add_paragraph()
            p.space_before = Pt(12)
            r1 = p.add_run(); r1.text = b_t; r1.font.bold = True; r1.font.size = Pt(14); r1.font.color.rgb = COLOR_TEXT
            r2 = p.add_run(); r2.text = n_t; r2.font.size = Pt(13); r2.font.color.rgb = COLOR_MUTED

        add_card(s, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.2))
        tb2 = s.shapes.add_textbox(Inches(7.05), Inches(1.9), Inches(5.2), Inches(4.7))
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        p = tf2.paragraphs[0]
        p.text = "Our Architectural Solution"
        p.font.size = Pt(19)
        p.font.bold = True
        p.font.color.rgb = COLOR_GREEN
        for b_t, n_t in [
            ("Deterministic Math Core: ", "Statistical Process Control (28-day DoW baseline) + exact 3-Factor Shapley decomposition with guaranteed $0.00 residual error."),
            ("3-Model AI Diagnostic Triad: ", "Internal diagnostic (Model 1) + Live macro sentinel (Model 2) + Prescriptive ROI trajectory simulator (Model 3)."),
            ("Governed Safe Action Engine: ", "Explicit abstention protocol under ambiguity ($<25\\%$ confidence delta) prescribing low-cost canary validation tests.")
        ]:
            p = tf2.add_paragraph()
            p.space_before = Pt(12)
            r1 = p.add_run(); r1.text = b_t; r1.font.bold = True; r1.font.size = Pt(14); r1.font.color.rgb = COLOR_TEXT
            r2 = p.add_run(); r2.text = n_t; r2.font.size = Pt(13); r2.font.color.rgb = COLOR_MUTED

        # --- Slide 2/9: Governed Semantic Contract ---
        s = prs.slides.add_slide(blank_layout)
        set_dark_bg(s)
        add_header(s, "Governed Semantic Contract & Multi-Source Ingestion")
        sources = [
            ("Daily ERP Sales Transactions", "Daily Grain (Structured SQL)", "Contains transaction orders, SKU categories, quantities, unit prices, and confidential COGS.", "1,420 rows  ·  24h SLA", Inches(0.8)),
            ("Hourly Web Analytics Stream", "Hourly Grain (Event Stream)", "Captures visitor sessions, funnel drop-offs, bounce rates, and checkout initiation events.", "696 hourly bins  ·  1h SLA", Inches(4.8)),
            ("Weekly Support & Jira Logs", "Weekly Grain (Semi-Structured)", "Customer friction tickets, warehouse error logs, and shipping delay incident taxonomy.", "18 incidents  ·  168h SLA", Inches(8.8))
        ]
        for title, grain, desc, sla, left_pos in sources:
            add_card(s, left_pos, Inches(1.7), Inches(3.7), Inches(3.5))
            tb = s.shapes.add_textbox(left_pos + Inches(0.2), Inches(1.9), Inches(3.3), Inches(3.1))
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = COLOR_ACCENT
            p_sub = tf.add_paragraph(); p_sub.text = grain; p_sub.font.size = Pt(12); p_sub.font.color.rgb = COLOR_GREEN; p_sub.space_before = Pt(4)
            p_desc = tf.add_paragraph(); p_desc.text = desc; p_desc.font.size = Pt(13); p_desc.font.color.rgb = COLOR_MUTED; p_desc.space_before = Pt(10)
            p_sla = tf.add_paragraph(); p_sla.text = sla; p_sla.font.size = Pt(12); p_sla.font.bold = True; p_sla.font.color.rgb = COLOR_TEXT; p_sla.space_before = Pt(12)

        add_card(s, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.4))
        tb_kpi = s.shapes.add_textbox(Inches(1.1), Inches(5.6), Inches(11.1), Inches(1.1))
        tf_kpi = tb_kpi.text_frame; tf_kpi.word_wrap = True
        p = tf_kpi.paragraphs[0]; p.text = "Governed 4-KPI Connected Equation (Strict Mathematical Lineage)"; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = COLOR_TEXT
        p_eq = tf_kpi.add_paragraph(); p_eq.text = "Gross Revenue (R) = Sessions (S) × Conversion Rate (CR) × Average Order Value (AOV)"; p_eq.font.size = Pt(16); p_eq.font.bold = True; p_eq.font.color.rgb = COLOR_ACCENT; p_eq.space_before = Pt(4)

        # --- Slide 3/9: Deterministic Math Core ---
        s = prs.slides.add_slide(blank_layout)
        set_dark_bg(s)
        add_header(s, "Deterministic Non-LLM Analytical Core")
        add_card(s, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.2))
        tb_spc = s.shapes.add_textbox(Inches(1.05), Inches(1.9), Inches(5.1), Inches(4.7))
        tf_spc = tb_spc.text_frame; tf_spc.word_wrap = True
        p = tf_spc.paragraphs[0]; p.text = "1. Statistical Process Control (SPC)"; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = COLOR_ACCENT
        for b_t, b_d in [
            ("Signal vs. Noise Filtering: ", "Applies 28-day rolling Day-of-Week seasonality index to prevent false alarms on regular weekend volume dips."),
            ("Dynamic Control Limits: ", "Computes UCL and LCL at $\\pm 2.5\\sigma$. Identifies Day 29 anomaly ($z = -20.75\\sigma$) as a true operational failure."),
            ("Zero LLM Tokens: ", "Executed purely in vectorized NumPy/Pandas in $<5\\text{ms}$ with zero API cost.")
        ]:
            p = tf_spc.add_paragraph(); p.space_before = Pt(12)
            r1 = p.add_run(); r1.text = b_t; r1.font.bold = True; r1.font.size = Pt(13); r1.font.color.rgb = COLOR_TEXT
            r2 = p.add_run(); r2.text = b_d; r2.font.size = Pt(12); r2.font.color.rgb = COLOR_MUTED

        add_card(s, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.2))
        tb_tree = s.shapes.add_textbox(Inches(7.05), Inches(1.9), Inches(5.2), Inches(4.7))
        tf_tree = tb_tree.text_frame; tf_tree.word_wrap = True
        p = tf_tree.paragraphs[0]; p.text = "2. Exact 3-Factor Shapley Decomposition"; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = COLOR_GREEN
        for b_t, b_d in [
            ("Exact Closed-Form Math: ", "Decomposes $-\\$65,600$ revenue drop across all factor permutation orderings (Sessions, CVR, AOV)."),
            ("Attribution Breakdown: ", "• Volume Impact: -$35,166 (-53.6%)\n• Conversion Rate: -$22,667 (-34.6%)\n• Basket Size (AOV): -$7,767 (-11.8%)"),
            ("Zero Residual Guarantee: ", "Sum of decomposed factors equals total revenue delta exactly: Residual $\\epsilon = \\$0.00$.")
        ]:
            p = tf_tree.add_paragraph(); p.space_before = Pt(12)
            r1 = p.add_run(); r1.text = b_t; r1.font.bold = True; r1.font.size = Pt(13); r1.font.color.rgb = COLOR_TEXT
            r2 = p.add_run(); r2.text = b_d; r2.font.size = Pt(12); r2.font.color.rgb = COLOR_MUTED

        # --- Slide 4/9: 3-Model AI Diagnostic Triad ---
        s = prs.slides.add_slide(blank_layout)
        set_dark_bg(s)
        add_header(s, "3-Model Cognitive AI Synthesis Triad")
        models = [
            ("Model 1: Diagnostic Engine", "Internal Operational Causes", "Analyzes internal ERP staging queues and Jira tickets (JIRA-4819). Isolates a 30% internal backlog at warehouse WH-WEST-01.", "30% Attribution Share", Inches(0.8), COLOR_PRIMARY),
            ("Model 2: Macro Sentinel", "External Live Market Feeds", "Ingests maritime port congestion feeds (MACRO-PORT-01). Quantifies a 70% external macro shock from West Coast labor slowdowns.", "70% Attribution Share", Inches(4.8), COLOR_AMBER),
            ("Model 3: Prescriptive Action", "ROI & Trajectory Simulator", "Blends Model 1 & 2 attributions, factors in dynamic executive constraints, and simulates 30/60/90-day recovery trajectories.", "4.2x Projected ROI", Inches(8.8), COLOR_GREEN)
        ]
        for title, subtitle, desc, metric, left_pos, border_c in models:
            add_card(s, left_pos, Inches(1.7), Inches(3.7), Inches(5.2), border_color=border_c)
            tb = s.shapes.add_textbox(left_pos + Inches(0.2), Inches(1.9), Inches(3.3), Inches(4.7))
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = COLOR_TEXT
            p_sub = tf.add_paragraph(); p_sub.text = subtitle; p_sub.font.size = Pt(12); p_sub.font.color.rgb = border_c; p_sub.space_before = Pt(4)
            p_desc = tf.add_paragraph(); p_desc.text = desc; p_desc.font.size = Pt(13); p_desc.font.color.rgb = COLOR_MUTED; p_desc.space_before = Pt(14)
            p_met = tf.add_paragraph(); p_met.text = metric; p_met.font.size = Pt(14); p_met.font.bold = True; p_met.font.color.rgb = border_c; p_met.space_before = Pt(20)

        # --- Slide 5/9: AI Safety & Explicit Abstention ---
        s = prs.slides.add_slide(blank_layout)
        set_dark_bg(s)
        add_header(s, "AI Safety: Explicit Abstention & Canary Tests (Scenario 2)")
        add_card(s, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.2))
        tb_ab1 = s.shapes.add_textbox(Inches(1.05), Inches(1.9), Inches(5.1), Inches(4.7))
        tf_ab1 = tb_ab1.text_frame; tf_ab1.word_wrap = True
        p = tf_ab1.paragraphs[0]; p.text = "When & Why We Abstain"; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = COLOR_AMBER
        for b_t, b_d in [
            ("The Capital Allocation Hazard: ", "Under ambiguous signals, forcing a single generative LLM prediction risks misallocating hundreds of thousands in capital."),
            ("Mathematical Abstention Gate: ", "Triggered automatically when the confidence margin between competing root causes is $<25\\%$ (observed: $16\\%$ delta)."),
            ("Dual Competing Hypotheses: ", "• Hypothesis 1 (60%): Payment Gateway 504 Timeout on iOS Checkout.\n• Hypothesis 2 (40%): Competitor Viral TikTok Flash Sale.")
        ]:
            p = tf_ab1.add_paragraph(); p.space_before = Pt(12)
            r1 = p.add_run(); r1.text = b_t; r1.font.bold = True; r1.font.size = Pt(13); r1.font.color.rgb = COLOR_TEXT
            r2 = p.add_run(); r2.text = b_d; r2.font.size = Pt(12); r2.font.color.rgb = COLOR_MUTED

        add_card(s, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.2))
        tb_ab2 = s.shapes.add_textbox(Inches(7.05), Inches(1.9), Inches(5.2), Inches(4.7))
        tf_ab2 = tb_ab2.text_frame; tf_ab2.word_wrap = True
        p = tf_ab2.paragraphs[0]; p.text = "Prescribed Low-Cost Discovery Tests"; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = COLOR_GREEN
        for title, body in [
            ("CANARY TEST 1: Synthetic iOS Checkout Probe", "Simulate 500 programmatic checkout requests through iOS gateway corridor.\nCost: $150  |  Duration: 2.0 hours  |  Decision Gate: Timeout Rate > 5%"),
            ("CANARY TEST 2: Competitive Ad & Social Scraping Run", "Ingest competitor TikTok and Instagram ad feeds across ASEAN.\nCost: $300  |  Duration: 4.0 hours  |  Decision Gate: Competitor Share of Voice Spike > 25%")
        ]:
            p = tf_ab2.add_paragraph(); p.space_before = Pt(14)
            r1 = p.add_run(); r1.text = title + "\n"; r1.font.bold = True; r1.font.size = Pt(13); r1.font.color.rgb = COLOR_ACCENT
            r2 = p.add_run(); r2.text = body; r2.font.size = Pt(12); r2.font.color.rgb = COLOR_MUTED

        # --- Slide 6/9: Persona Prescriptive Storytelling & RBAC ---
        s = prs.slides.add_slide(blank_layout)
        set_dark_bg(s)
        add_header(s, "Persona Prescriptive Storytelling & RBAC Security")
        add_card(s, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.2))
        tb_ex = s.shapes.add_textbox(Inches(1.05), Inches(1.9), Inches(5.1), Inches(4.7))
        tf_ex = tb_ex.text_frame; tf_ex.word_wrap = True
        p = tf_ex.paragraphs[0]; p.text = "Executive / VP Persona"; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = COLOR_PRIMARY
        for b_t, b_d in [
            ("Tailored Decision Lens: ", "High-level strategic trade-offs, financial risk exposure, capital ROI, and supply chain corridor reallocation."),
            ("Prescribed Interventions: ", "• Authorize 35% maritime freight reroute to Seattle port corridors.\n• Authorize $45,000 expedited drayage budget to clear high-margin inventory."),
            ("Unmasked Financial Data: ", "Full numeric transparency over unit COGS, gross margin amounts, and margin percentages.")
        ]:
            p = tf_ex.add_paragraph(); p.space_before = Pt(12)
            r1 = p.add_run(); r1.text = b_t; r1.font.bold = True; r1.font.size = Pt(13); r1.font.color.rgb = COLOR_TEXT
            r2 = p.add_run(); r2.text = b_d; r2.font.size = Pt(12); r2.font.color.rgb = COLOR_MUTED

        add_card(s, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.2))
        tb_op = s.shapes.add_textbox(Inches(7.05), Inches(1.9), Inches(5.2), Inches(4.7))
        tf_op = tb_op.text_frame; tf_op.word_wrap = True
        p = tf_op.paragraphs[0]; p.text = "Operations / Category Manager"; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = COLOR_ACCENT
        for b_t, b_d in [
            ("Tailored Decision Lens: ", "Tactical SOPs, warehouse triage playbooks, system container codes, and ticket resolution."),
            ("Prescribed Interventions: ", "• Execute automated inventory reallocation for SKU-ELEC-409.\n• Triage Jira incident JIRA-4819 at Long Beach Berth 4.\n• Trigger customer delivery ETA updates to enterprise accounts."),
            ("Enforced RBAC Masking: ", "Unit COGS and Gross Margins are strictly redacted with [RESTRICTED-EXEC].")
        ]:
            p = tf_op.add_paragraph(); p.space_before = Pt(12)
            r1 = p.add_run(); r1.text = b_t; r1.font.bold = True; r1.font.size = Pt(13); r1.font.color.rgb = COLOR_TEXT
            r2 = p.add_run(); r2.text = b_d; r2.font.size = Pt(12); r2.font.color.rgb = COLOR_MUTED

        # --- Slide 7/9: Continuous Learning Loop & Telemetry ---
        s = prs.slides.add_slide(blank_layout)
        set_dark_bg(s)
        add_header(s, "Human-in-the-Loop Feedback & Telemetry Audit")
        add_card(s, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.2))
        tb_fb = s.shapes.add_textbox(Inches(1.05), Inches(1.9), Inches(5.1), Inches(4.7))
        tf_fb = tb_fb.text_frame; tf_fb.word_wrap = True
        p = tf_fb.paragraphs[0]; p.text = "Continuous Learning Loop"; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = COLOR_ACCENT
        for b_t, b_d in [
            ("Governed Human Corrections: ", "Analysts submit star ratings, categorical tags, and operational driver corrections."),
            ("Cognitive Interpretation: ", "AI classifies feedback type, identifies target layer, and generates learning signals."),
            ("Bayesian Calibration Weights: ", "Updates operational driver ranking weights over time."),
            ("Quantitative Lock Principle: ", "Closed-form formulas, variance math, and SPC thresholds remain strictly immutable.")
        ]:
            p = tf_fb.add_paragraph(); p.space_before = Pt(10)
            r1 = p.add_run(); r1.text = b_t; r1.font.bold = True; r1.font.size = Pt(13); r1.font.color.rgb = COLOR_TEXT
            r2 = p.add_run(); r2.text = b_d; r2.font.size = Pt(12); r2.font.color.rgb = COLOR_MUTED

        add_card(s, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.2))
        tb_tel = s.shapes.add_textbox(Inches(7.05), Inches(1.9), Inches(5.2), Inches(4.7))
        tf_tel = tb_tel.text_frame; tf_tel.word_wrap = True
        p = tf_tel.paragraphs[0]; p.text = "Runtime Telemetry & Cost Accounting"; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = COLOR_GREEN
        for b_t, b_d in [
            ("Total Pipeline Latency: ", "< 2,000 ms (Ingest 1.5ms | SPC & Tree Math 4.2ms | AI Synthesis 1,600ms)."),
            ("Deterministic Math Tokens: ", "0 Tokens consumed (100% Non-LLM Closed-Form Python/NumPy)."),
            ("AI Narrative Tokens: ", "630 Tokens (450 prompt + 180 completion tokens)."),
            ("Cost Transparency: ", "$0.001035 estimated cost per execution (with $0.00 offline fallback).")
        ]:
            p = tf_tel.add_paragraph(); p.space_before = Pt(10)
            r1 = p.add_run(); r1.text = b_t; r1.font.bold = True; r1.font.size = Pt(13); r1.font.color.rgb = COLOR_TEXT
            r2 = p.add_run(); r2.text = b_d; r2.font.size = Pt(12); r2.font.color.rgb = COLOR_MUTED

        # --- Slide 8/9: Complete 4-Scenario Matrix ---
        s = prs.slides.add_slide(blank_layout)
        set_dark_bg(s)
        add_header(s, "Complete 4-Scenario Validation Matrix")
        scenarios = [
            ("Scenario 1: Multi-Factor Disruption", "Compound Drivers", "70% External West Coast port strike + 30% local warehouse backlog. 4.2x projected ROI recovery.", Inches(0.8), Inches(1.7)),
            ("Scenario 2: Low-Confidence Ambiguity", "Explicit Abstention", "Conflicting evidence (60% vs 40%). Engine abstains and prescribes $150 canary validation experiments.", Inches(6.8), Inches(1.7)),
            ("Scenario 3: Sparse History / Launch", "Cold Start Prior", "New SKU launch with N=7 days history. Blends Bayesian category benchmark ($5,000 base) and staged pilots.", Inches(0.8), Inches(4.4)),
            ("Scenario 4: Role-Based Security", "RBAC Entitlements", "Enforces column/row masking on unit COGS and margins for Analysts while showing system IDs.", Inches(6.8), Inches(4.4))
        ]
        for title, subtitle, desc, left_pos, top_pos in scenarios:
            add_card(s, left_pos, top_pos, Inches(5.7), Inches(2.4))
            tb = s.shapes.add_textbox(left_pos + Inches(0.2), top_pos + Inches(0.15), Inches(5.3), Inches(2.1))
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = COLOR_ACCENT
            p_sub = tf.add_paragraph(); p_sub.text = subtitle; p_sub.font.size = Pt(12); p_sub.font.bold = True; p_sub.font.color.rgb = COLOR_GREEN; p_sub.space_before = Pt(2)
            p_desc = tf.add_paragraph(); p_desc.text = desc; p_desc.font.size = Pt(13); p_desc.font.color.rgb = COLOR_MUTED; p_desc.space_before = Pt(6)

        # --- Slide 9/9: Summary & Architectural Differentiators ---
        s = prs.slides.add_slide(blank_layout)
        set_dark_bg(s)
        add_header(s, "Summary: Production-Grade KPI Intelligence Engine")
        add_card(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(5.2))
        tb_end = s.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(11.0), Inches(4.6))
        tf_end = tb_end.text_frame; tf_end.word_wrap = True
        p = tf_end.paragraphs[0]; p.text = "Key Architectural Differentiators"; p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = COLOR_TEXT
        for b_t, b_d in [
            ("1. Mathematical Rigor First: ", "Zero reliance on LLM calculation. SPC anomaly detection and Shapley causal attribution are 100% deterministic with $0.00 residual."),
            ("2. Multi-Model Context Synthesis: ", "Combines private ERP/Jira operational context with live macro market shocks into actionable, persona-tailored playbooks."),
            ("3. Enterprise Guardrails: ", "Explicit abstention under ambiguity, dynamic RBAC masking, and governed human feedback calibration."),
            ("4. Verified & Ready: ", "93 / 93 automated tests passed. 8 / 8 scenario-persona combinations validated. Deployed live on Streamlit Cloud.")
        ]:
            p = tf_end.add_paragraph(); p.space_before = Pt(14)
            r1 = p.add_run(); r1.text = b_t; r1.font.bold = True; r1.font.size = Pt(15); r1.font.color.rgb = COLOR_ACCENT
            r2 = p.add_run(); r2.text = b_d; r2.font.size = Pt(14); r2.font.color.rgb = COLOR_MUTED

    # Build the middle slides
    # Currently prs has 3 slides: [0: Cover, 1: Team, 2: OG Thank You]
    # We will temporarily remove slide 2 (OG Thank You) element from slide list, append the 9 middle slides, and re-append OG Thank You element at the end!
    og_thank_you_elem = prs.slides._sldIdLst[2]
    prs.slides._sldIdLst.remove(og_thank_you_elem)

    # Add middle slides
    build_middle_slides()

    # Re-append OG Thank You at the very end
    prs.slides._sldIdLst.append(og_thank_you_elem)

    for target in [output_path, "Warriors_Innovation_Challenge_2026.pptx", "BusinessIntelligence_Executive_Deck.pptx"]:
        try:
            prs.save(target)
            print(f"Merged deck saved successfully with {len(prs.slides)} slides to: {target}")
        except PermissionError:
            print(f"Note: {target} is currently open in PowerPoint. Saved to alternate file.")

if __name__ == "__main__":
    b_path = "Warriors_ 3-Model Triad_backup.pptx"
    o_path = "Warriors_ 3-Model Triad.pptx"
    build_merged_warriors_deck(b_path, o_path)
